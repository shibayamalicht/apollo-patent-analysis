# -*- coding: utf-8 -*-
"""APOLLO CAPCOM スライド生成ヘルパーライブラリ — v6 "Mission Deck" デザイン。

PPTX 生成スクリプトから `from apollo_slides import *` で使う。

デザイン原則（v6 "Mission Deck"・ユーザー承認済みサンプル準拠）:
  - 2層サーフェス: オフホワイト地 (BG) に白カード (CARD, 角丸, ヘアライン枠)
  - インク主体のタイポグラフィ: 数値・本文は Navy インク。色は「意味」にだけ使う
  - カテゴリ4色 (青/teal/紫/琥珀) はアクセントバー・チップ・バンド地に載せ、
    文字には載せない（白文字を載せるのは濃色変種 CHIP / DEEP のみ）
  - 帯 (バンド) は淡色ティント地 + 彩色左バー + 白抜きチップ
  - Noto Sans JP 多段ウェイト + 全ランに lang=ja-JP（中華風グリフの構造的防止）
  - 見出し・数値は Noto Sans JP + bold フラグ（Black は章扉ゴースト数字のみ）
  - 裏表紙に "Thank You" は置かない

呼び出し側で必要な前提:
    from pptx import Presentation
    prs = Presentation(TEMPLATE_OR_OWN_TEMPLATE)
    blank = prs.slide_layouts[6]
    SNAP = "<セッションフォルダ>/snapshots"   # 必要なら本モジュールの SNAP を上書き
"""

from pptx import Presentation  # noqa: F401  （呼び出しスクリプトが prs を作る際に使う）
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from lxml import etree
import os
import re
import math

# =============================================================================
# デザイントークン（v6 "Mission Deck"）
# =============================================================================
BG        = RGBColor(0xF4, 0xF6, 0xF9)   # スライド地（クールオフホワイト）
CARD      = RGBColor(0xFF, 0xFF, 0xFF)   # カード地（白）
CARD_LINE = RGBColor(0xE3, 0xE8, 0xF1)   # カード枠・区切り線（ヘアライン）
INK       = RGBColor(0x1B, 0x2A, 0x4A)   # 主文字色（Navy インク）
DEEP      = RGBColor(0x15, 0x22, 0x38)   # 濃紺（表紙・裏表紙・結論ピル・表ヘッダー）
SUB       = RGBColor(0x4F, 0x5B, 0x70)   # 副文字色
MUTED     = RGBColor(0x87, 0x92, 0xA3)   # 出所・フッター・キャプション
GHOST     = RGBColor(0xE9, 0xEE, 0xF6)   # 章扉ゴースト数字

ACCENT    = RGBColor(0x2A, 0x78, 0xD6)   # ブランドアクセント（青 = カテゴリ1と同一）
CAT = {                                   # カテゴリ4色（彩色バー用）
    "blue":   RGBColor(0x2A, 0x78, 0xD6),
    "teal":   RGBColor(0x1B, 0xAF, 0x7A),
    "violet": RGBColor(0x4A, 0x3A, 0xA7),
    "amber":  RGBColor(0xED, 0xA1, 0x00),
}
TINT = {                                  # バンド地・カードヘッダ用の淡色ティント
    "blue":   RGBColor(0xE9, 0xF1, 0xFB),
    "teal":   RGBColor(0xE5, 0xF6, 0xEF),
    "violet": RGBColor(0xED, 0xEA, 0xF7),
    "amber":  RGBColor(0xFC, 0xF3, 0xDD),
}
CHIP = {                                  # 白文字を載せる濃色変種（コントラスト確保）
    "blue":   RGBColor(0x1F, 0x5F, 0xB0),
    "teal":   RGBColor(0x12, 0x80, 0x5A),
    "violet": RGBColor(0x4A, 0x3A, 0xA7),
    "amber":  RGBColor(0xB2, 0x6A, 0x00),
}
GOOD_TX   = RGBColor(0x1E, 0x7B, 0x34)   # 判定文字色: 支持（緑）
PART_TX   = RGBColor(0xB2, 0x6A, 0x00)   # 判定文字色: 部分支持（琥珀）
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SKY       = RGBColor(0x8F, 0xB8, 0xE8)   # 濃紺地の上のアクセント文字（明青）
SUB_DK    = RGBColor(0xB9, 0xC6, 0xDC)   # 濃紺地の上の副文字
MUT_DK    = RGBColor(0x76, 0x87, 0xA5)   # 濃紺地の上の弱文字
RED_ACCENT = RGBColor(0xD6, 0x45, 0x45)  # 警告・却下・優先度「高」

CAT_ORDER = ("blue", "teal", "violet", "amber")  # カテゴリ固定割当の順序

# --- 後方互換エイリアス（旧定数名。値は新体系に合わせて再定義） -----------------
NAVY = INK                                # 旧: タイトル・強調テキスト
BLUE = CHIP["blue"]                       # 旧: セクションヘッダー背景
DARK_GRAY = INK                           # 旧: 本文テキスト（インク主体に統一）
MEDIUM_GRAY = SUB                         # 旧: 補足テキスト・キャプション
LIGHT_GRAY = RGBColor(0xED, 0xF1, 0xF7)   # 旧: 枠背景（残存互換用の淡グレー）
BORDER_GRAY = CARD_LINE                   # 旧: 罫線・区切り線
KEY_MSG_BG = TINT["blue"]                 # 旧: 強調ボックス背景
GREEN_ACCENT = CAT["teal"]                # 旧: ポジティブ指標
AMBER = CAT["amber"]                      # 旧: 注意指標
GHOST_NAVY = GHOST                        # 旧: セクションゴースト番号

# ボトムバー定数（後方互換のため残置。v6 ではフッターはヘアライン線のみ）
BOTTOM_BAR_HEIGHT = 4  # px（旧仕様の名残）
BOTTOM_BAR_Y = 6.92    # Inches — フッター線の直上

# バー太さ（EMU）
BAR_W_LEFT = Emu(45720)   # カード左のカテゴリ色バー 3.6pt
BAR_H_TOP = Emu(41148)    # カード上のカテゴリ色バー 3.2pt
HAIRLINE = Emu(9525)      # 0.75pt
UNDERLINE_H = Emu(28575)  # タイトル短下線 2.25pt
UNDERLINE_W = 0.55        # タイトル短下線の幅 (inch)
RADIUS_IN = 0.07          # カード角丸の絶対半径 (inch)

# =============================================================================
# フォント設定（Noto Sans JP・多段ウェイト）
# =============================================================================
FONT_FAMILY = "Noto Sans JP"            # 既定（Regular）
# 役割別ウェイト。Noto Sans JP の静的ウェイト名（未インストール時は近いウェイトに自動フォールバック）。
WEIGHT_FAMILY = {
    "light":    "Noto Sans JP Light",     # キャプション・出典・フッター（控えめ）
    "regular":  "Noto Sans JP",           # 本文・注釈
    "medium":   "Noto Sans JP Medium",    # 小見出し・リード文・チップ
    "semibold": "Noto Sans JP SemiBold",  # 強調本文・カード見出し
    "bold":     "Noto Sans JP",           # 見出し・数値は bold フラグで（Black より一段上品）
    "black":    "Noto Sans JP Black",     # 章扉ゴースト数字のみ
}
JA_FONT = FONT_FAMILY                    # 後方互換
LATIN_FONT = FONT_FAMILY

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# =============================================================================
# パス定数
# =============================================================================
# テンプレートPPTX（呼び出しスクリプトの位置を基準に解決する想定の相対パス）。
# 本モジュールでは import 時に Presentation を生成しない（生成は呼び出し側の責務）。
TEMPLATE = os.path.join(os.path.dirname(__file__), "../../capcom_schema/templates/apollo_template.pptx")
# スナップショット画像フォルダ。fit_image / add_chart_text_slide 等が相対パス画像を
# 解決する際の基点。CAPCOM セッション内の相対パス "snapshots" を既定にする。
# 呼び出しスクリプト側で実フォルダ（例: os.path.join(session_dir, "snapshots")）に
# 上書きしてよい。
SNAP = "snapshots"


# =============================================================================
# Section 2: コアユーティリティ — フォント・禁則ヘルパー
# =============================================================================
def _apply_font(run, weight=None):
    """runにデュアルフォント（欧文 + 日本語）+ 言語 + ウェイトを設定する。

    weight: None/"regular"/"medium"/"semibold"/"bold"/"black"/"light"。
      - None: 既定（Noto Sans JP Regular）。呼び出し側の bold フラグはそのまま尊重。
      - 名前付きウェイト（light/medium/semibold/black 等）: その専用ファミリを使い、
        bold フラグは使わない（名前付きウェイトと bold の二重指定を避ける）。
      - "bold": Regular ファミリ + bold フラグ（Bold 相当）。
    """
    fam = WEIGHT_FAMILY.get(weight, FONT_FAMILY) if weight else FONT_FAMILY
    run.font.name = fam
    if weight in ("light", "regular", "medium", "semibold", "black"):
        run.font.bold = False            # 名前付きウェイトでは bold を使わない
    elif weight == "bold":
        run.font.bold = True
    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'ja-JP')
    rPr.set('altLang', 'en-US')
    ea = rPr.find(f'{{{A_NS}}}ea')
    if ea is None:
        ea = etree.SubElement(rPr, f'{{{A_NS}}}ea')
    ea.set('typeface', fam)


def _apply_kinsoku(paragraph):
    """段落に日本語禁則処理を設定する"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('eaLnBrk', '1')
    pPr.set('hangingPunct', '1')


def _set_spacing(run, spc):
    """runに字間（spc: 1/100pt 単位）を設定する（アイブロウ・チップ用）"""
    try:
        run._r.get_or_add_rPr().set('spc', str(spc))
    except Exception:
        pass


# =============================================================================
# 低レベル図形ヘルパー（2層サーフェス・カード・バー）
# =============================================================================
def _in(v):
    """float(インチ) → Inches。Length 型（Emu/Pt 等）はそのまま返す"""
    return Inches(v) if isinstance(v, (int, float)) else v


def _shadow_off(shape):
    """図形の継承シャドウを無効化する"""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _rect(slide, x, y, w, h, fill, line=None, line_w=HAIRLINE,
          round_=False, radius_in=RADIUS_IN):
    """矩形/角丸矩形。fill=None で塗りなし、line=None で枠線なし。影は常に無効"""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        _in(x), _in(y), _in(w), _in(h))
    if round_:
        try:  # 角丸を絶対半径 radius_in (inch) に固定（大きさで半径が暴れない）
            shp.adjustments[0] = max(0.02, min(0.5, radius_in / min(shp.width.inches, shp.height.inches)))
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    _shadow_off(shp)
    return shp


def _set_bg(slide, color=BG):
    """スライド背景色を設定する（コンテンツ系=BG、表紙/裏表紙=DEEP）"""
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _card(slide, x, y, w, h, fill=CARD, line=CARD_LINE):
    """共通カード部品 — 白地・角丸(絶対半径0.07in)・ヘアライン枠・影なし"""
    return _rect(slide, x, y, w, h, fill, line, round_=True)


def _bar_left(slide, x, y, h, color):
    """カード左のカテゴリ色バー（3.6pt）"""
    return _rect(slide, x, y, BAR_W_LEFT, h, color)


def _bar_top(slide, x, y, w, color):
    """カード上のカテゴリ色バー（3.2pt・左右0.1inインセットで角丸を避ける）"""
    return _rect(slide, x + 0.10, y, w - 0.20, BAR_H_TOP, color)


_RGB_TO_KEY = {}
for _k in CAT:
    _RGB_TO_KEY[CAT[_k]] = _k
    _RGB_TO_KEY[CHIP[_k]] = _k
_RGB_TO_KEY[INK] = "blue"
_RGB_TO_KEY[DEEP] = "blue"
_RGB_TO_KEY[RED_ACCENT] = "amber"
_RGB_TO_KEY[SUB] = "violet"


def _cat_key(color, idx=0):
    """色指定（"blue" 等のキー / RGBColor / None）をカテゴリキーに解決する。

    None はカテゴリ順の固定割当（CAT_ORDER[idx % 4]）。旧 API から渡される
    RGBColor は最も近い意味のカテゴリへマップする。
    """
    if color is None:
        return CAT_ORDER[idx % 4]
    if isinstance(color, str):
        return color if color in CAT else CAT_ORDER[idx % 4]
    return _RGB_TO_KEY.get(color, CAT_ORDER[idx % 4])


def _chip_pill(slide, x, y, text, color_key="blue", w=1.15, h=0.30, size=9.5):
    """白抜きチップ（CHIP 濃色地 + 白文字 + 字間広め）"""
    c = _rect(slide, x, y, w, h, CHIP.get(color_key, CHIP["blue"]),
              round_=True, radius_in=0.05)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = WHITE
    _apply_font(run, "medium")
    _set_spacing(run, 60)
    return c


def _strip_table_borders(table):
    """テーブルの全セル罫線を透明化する（罫線レス化）"""
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ('lnL', 'lnR', 'lnT', 'lnB'):
                el = tcPr.find(f'{{{A_NS}}}{tag}')
                if el is not None:
                    tcPr.remove(el)
            for i, tag in enumerate(('lnL', 'lnR', 'lnT', 'lnB')):
                ln = tcPr.makeelement(f'{{{A_NS}}}{tag}', {'w': '0'})
                etree.SubElement(ln, f'{{{A_NS}}}noFill')
                tcPr.insert(i, ln)


def _plain_table_style(table):
    """python-pptx 既定のテーブルスタイル装飾（先頭行強調・縞）を無効化する"""
    try:
        table.first_row = False
        table.horz_banding = False
    except Exception:
        pass


# =============================================================================
# テキストエンジン
# =============================================================================
def add_rich_runs(paragraph, text, base_size=Pt(14), base_color=DARK_GRAY,
                  bold_color=None, force_bold=False, line_spacing=1.4, weight=None):
    """**太字**マーカー解析 + デュアルフォント + 禁則 + 行間 + ウェイト

    weight を指定すると全 run にそのウェイト（例 "bold"）を適用する。
    weight=None の場合は **マーカー** 部のみ Bold、他は force_bold に従う（従来動作）。
    """
    paragraph.clear()
    bold_color = bold_color or base_color
    _apply_kinsoku(paragraph)
    paragraph.line_spacing = line_spacing

    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.size = base_size
            run.font.bold = True
            run.font.color.rgb = bold_color
        else:
            run = paragraph.add_run()
            run.text = part
            run.font.size = base_size
            run.font.bold = force_bold
            run.font.color.rgb = bold_color if force_bold else base_color
        _apply_font(run, weight)


def set_text(p, text, size, color, bold=False, line_spacing=None, weight=None):
    """シンプルテキスト設定（デュアルフォント + 禁則 + ウェイト付き）

    weight を指定すると名前付きウェイト（light/medium/bold 等）を適用する
    （bold 引数より優先）。未指定なら bold フラグで太字制御（従来動作）。
    """
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    _apply_font(run, weight)
    _apply_kinsoku(p)
    if line_spacing:
        p.line_spacing = line_spacing


# =============================================================================
# スライドタイトル（結論型 + 短下線）
# =============================================================================
def add_title_shape(slide, text, x=0.5, y=0.15, w=12.3, eyebrow=None):
    """スライドタイトル（結論見出し・INK・bold + 左端 0.55in の短下線）。

    結論を 1 文で言い切る（「～」副題は使わない。必要なら全角ダッシュ「—」か
    句点で短く 2 文に分ける）。数値を必ず含める。Noto Sans JP + bold フラグを使用
    （Black は章扉ゴースト数字専用）。下線は全幅ではなく左端 0.55in の
    アクセント短下線（2.25pt ACCENT）。

    eyebrow（任意）: タイトル直上に小さな「アイブロウ」（章/モジュール名。例
        "NEBULA / 環境分析"）を添える。Noto Sans JP Medium・10pt・ミュート色・
        字間を広げて、編集的な見出し階層（アイブロウ→主張見出し→リード文→
        根拠→締め文・§0.9-A0）を再現する。明朝/等幅は使わずゴシックで統一。
    Returns:
        float: タイトル下端のy座標（サブメッセージの配置基準）
    """
    # アイブロウ（任意）— タイトルの上に章/モジュール名の小ラベルを置く
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.26))
        eb.text_frame.word_wrap = True
        ep = eb.text_frame.paragraphs[0]
        set_text(ep, str(eyebrow), Pt(10), MUTED, weight="medium")
        for _r in ep.runs:  # 字間を少し広げてエディトリアルなアイブロウに（spc=1/100pt）
            _set_spacing(_r, 180)
        y = y + 0.30  # タイトルをアイブロウ分だけ下げる

    text_len = len(text)
    if text_len <= 30:
        font_size = Pt(22)
        box_h = 0.60
    elif text_len <= 50:
        font_size = Pt(20)
        box_h = 0.70
    else:
        font_size = Pt(18)
        box_h = 0.85

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    add_rich_runs(p, text, base_size=font_size, base_color=INK,
                  bold_color=INK, line_spacing=1.25, weight="bold")

    # 短下線 — 左端 0.55in・2.25pt・ACCENT
    line_y = y + box_h + 0.03
    _rect(slide, x, line_y, UNDERLINE_W, UNDERLINE_H, ACCENT)

    return line_y + 0.09  # サブメッセージの開始y座標を返す


# =============================================================================
# リード文（サブメッセージ）
# =============================================================================
def add_sub_message(slide, message, x=0.5, y=None, w=12.3):
    """リード文（タイトル直下のプレーンな導入文・12.5pt Medium・SUB色）。

    v6 で ■マーカー + 強調ボックス（KEY_MSG_BG 箱）は廃止した。
    ⚠️ 本関数は「■」を自動付与しない（旧仕様の■自動付与は廃止）。
       `message` に先頭「■」を付けないこと。万一付いていても先頭の■は除去される。
    **強調** マーカーは INK 色の SemiBold で表現される。
    Args:
        y: 開始y座標。Noneの場合はadd_title_shapeの戻り値を使うこと。
    Returns:
        float: リード文下端のy座標 + マージン（コンテンツ開始位置）
    """
    # 防御: 旧仕様の名残で先頭に「■」が付いていても取り除く
    message = re.sub(r'^[\s　]*[■▪◾]\s*', '', message)
    if y is None:
        y = 1.00
    # 行数見積り（12.5pt・幅12.3in で1行あたり約65文字）
    est_lines = max(1, -(-len(message) // 62))
    box_h = 0.32 * est_lines + 0.04

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(box_h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    parts = re.split(r'(\*\*.*?\*\*)', message)
    for part in parts:
        if not part:
            continue
        run = p.add_run()
        run.font.size = Pt(12.5)
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.color.rgb = INK
            _apply_font(run, "semibold")           # 強調は SemiBold・インク色
        else:
            run.text = part
            run.font.color.rgb = SUB
            _apply_font(run, "medium")             # リード文は Medium・SUB色
    _apply_kinsoku(p)
    p.line_spacing = 1.35

    return y + box_h + 0.14


# =============================================================================
# フッター（全コンテンツスライド必須）
# =============================================================================
def add_bottom_bar_and_footer(slide, page_num=None, report_title=None, date_str=None):
    """全コンテンツスライド共通フッター — ヘアライン区切り線 + 左情報 + 右ページ番号。

    v6 で NAVY 太バーは廃止し、ヘアライン線のみの軽いフッターに刷新した。
    左: "APOLLO｜{レポートタイトル}｜{日付}"（report_title/date_str 未指定分は省略）
    右: ページ番号（ゼロ埋め2桁）
    タイトルスライド・セクションスライド・クロージングスライドでは呼ばない。
    """
    # ヘアライン区切り線
    _rect(slide, 0.5, 6.95, 12.33, HAIRLINE, CARD_LINE)

    # 左: "APOLLO｜タイトル｜日付"
    left_parts = ["APOLLO"]
    if report_title:
        left_parts.append(str(report_title))
    if date_str:
        left_parts.append(str(date_str))
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.03), Inches(9.0), Inches(0.25))
    set_text(txBox.text_frame.paragraphs[0], "｜".join(left_parts), Pt(8), MUTED, weight="light")

    # 右: ページ番号
    if page_num is not None:
        page_text = f"{page_num:02d}" if isinstance(page_num, int) else str(page_num)
        txBox2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.03), Inches(1.33), Inches(0.25))
        p2 = txBox2.text_frame.paragraphs[0]
        set_text(p2, page_text, Pt(9), MUTED, weight="light")
        p2.alignment = PP_ALIGN.RIGHT


# =============================================================================
# 画像・データソース・注釈
# =============================================================================
def fit_image(slide, image_path, max_x, max_y, max_w, max_h):
    """画像をアスペクト比保持で指定領域内に中央配置"""
    if not os.path.exists(image_path):
        return None
    img = Image.open(image_path)
    img_w, img_h = img.size
    ratio = img_h / img_w
    if max_w * ratio <= max_h:
        use_w, use_h = max_w, max_w * ratio
    else:
        use_h, use_w = max_h, max_h / ratio
    left = max_x + (max_w - use_w) / 2
    top = max_y + (max_h - use_h) / 2
    pic = slide.shapes.add_picture(
        image_path, Inches(left), Inches(top),
        width=Inches(use_w), height=Inches(use_h)
    )
    img.close()
    return pic


def _card_image(slide, image_path, x, y, w, h, pad=0.12):
    """白カードの上に画像をアスペクト維持で配置する（カード + 画像の複合部品）"""
    _card(slide, x, y, w, h)
    return fit_image(slide, image_path, max_x=x + pad, max_y=y + pad,
                     max_w=w - pad * 2, max_h=h - pad * 2)


def add_source_label(slide, source_text, x=0.5, y=6.55, w=12.3):
    """（出所）ラベル"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    set_text(p, f"（出所）{source_text}", Pt(9), MUTED, weight="light")


def add_annotation_block(slide, bullets, x, y, w, h, font_size=14,
                         has_border=False, bg_color=None):
    """テキスト注釈ブロック（チャート横の分析テキスト）

    ■マーカー付き箇条書きでチャートを補足する。
    各bullet: 最大2行、14pt。全体で3-5項目を推奨。
    """
    if bg_color or has_border:
        box = _rect(slide, x, y, w, h,
                    bg_color if bg_color else None,
                    CARD_LINE if has_border else None,
                    round_=True)
        if bg_color is None:
            box.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.08),
        Inches(w - 0.24), Inches(h - 0.16)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        marker = p.add_run()
        marker.text = "■ "
        marker.font.size = Pt(font_size)
        marker.font.color.rgb = ACCENT
        _apply_font(marker)
        parts = re.split(r'(\*\*.*?\*\*)', item)
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
                run.font.color.rgb = INK
            else:
                run.text = part
                run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(font_size)
            _apply_font(run)
        _apply_kinsoku(p)
        p.line_spacing = 1.5


def add_chart_label(slide, text, x, y, w=3.0, size=14, color=INK):
    """チャート小見出し（グラフ上の分類ラベル）"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    set_text(txBox.text_frame.paragraphs[0], text, Pt(size), color, weight="semibold")


def _add_line(slide, x0, y0, x1, y1, color, weight_pt=1.5):
    """2点間を細い回転矩形で結ぶ（コネクタ p:cxnSp はファイル破損を起こすため使わない）。

    座標は Inches 値。線色・太さは元のコネクタの見た目を踏襲する。
    """
    dx, dy = (x1 - x0), (y1 - y0)
    length = math.hypot(dx, dy)
    angle = math.degrees(math.atan2(dy, dx))
    cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
    thick = weight_pt / 72.0  # pt → inch
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - length / 2), Inches(cy - thick / 2),
        Inches(length), Inches(thick)
    )
    ln.rotation = angle
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def add_chart_callout(slide, text, x, y, w=2.5,
                      arrow_to_x=None, arrow_to_y=None,
                      bg_color=None, font_size=12, border_color=None):
    """チャート上の吹き出し注釈（画像の上にオーバーレイ配置）"""
    bg_color = bg_color or WHITE
    border_color = border_color or ACCENT

    chars_per_line = int(w * 7)
    num_lines = max(1, -(-len(text) // chars_per_line))
    h = 0.15 + num_lines * 0.28

    box = _rect(slide, x, y, w, h, bg_color, border_color, line_w=Emu(12700),
                round_=True, radius_in=0.05)

    txBox = slide.shapes.add_textbox(
        Inches(x + 0.08), Inches(y + 0.04), Inches(w - 0.16), Inches(h - 0.08)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(tf.paragraphs[0], text, base_size=Pt(font_size),
                  base_color=INK, bold_color=INK, line_spacing=1.3)

    if arrow_to_x is not None and arrow_to_y is not None:
        cx = x + w / 2
        cy = y + h / 2
        # コネクタ（p:cxnSp）はファイル破損を起こすため、回転矩形で引き出し線を描く
        _add_line(slide, cx, cy, arrow_to_x, arrow_to_y, border_color, weight_pt=1.0)

    return box


def add_highlight_circle(slide, x, y, w=0.5, h=0.5, color=None):
    """チャート上のハイライト丸囲み"""
    color = color or RED_ACCENT
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    oval.fill.background()
    oval.line.color.rgb = color
    oval.line.width = Emu(19050)
    _shadow_off(oval)
    return oval


# =============================================================================
# 新部品: コールアウトボックス / 注釈カードレール
# =============================================================================
def add_callout_box(slide, x, y, w, h, title, body, kind='caution'):
    """コールアウトボックス部品（注意喚起 / キーポイント）。

    kind='caution'  : 琥珀ティント地 + 琥珀左バー + PART_TX 見出し（注意・限界の明示）
    kind='keypoint' : 青ティント地 + 青左バー + INK 見出し（決め手・キーポイント）
    """
    ck = "amber" if kind == 'caution' else "blue"
    title_color = PART_TX if kind == 'caution' else INK
    _rect(slide, x, y, w, h, TINT[ck], round_=True)
    _bar_left(slide, x, y, h, CAT[ck])
    txT = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12),
                                   Inches(w - 0.4), Inches(0.3))
    set_text(txT.text_frame.paragraphs[0], title, Pt(11), title_color, weight="semibold")
    txB = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.46),
                                   Inches(w - 0.4), Inches(max(0.2, h - 0.58)))
    tf = txB.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_rich_runs(tf.paragraphs[0], body, base_size=Pt(9.5),
                  base_color=INK, bold_color=INK, line_spacing=1.35)


def add_annotation_cards(slide, x, y, w, items, card_h=None):
    """注釈カードレール部品 — 白カード + カテゴリ色左バーの小カードを縦に並べる。

    items: [{"title":"図の見方", "body":"…1-2文", "color":"blue"}, ...]（2〜4推奨）
    color はカテゴリキー（"blue"/"teal"/"violet"/"amber"）。省略時はカテゴリ順に固定割当。
    card_h 省略時は下端 6.4in までに収まるよう自動計算する。
    Returns:
        float: 最終カード下端のy座標
    """
    n = max(1, len(items))
    gap = 0.11
    if card_h is None:
        card_h = min(1.5, max(0.7, (6.4 - y - gap * (n - 1)) / n))
    for i, item in enumerate(items):
        cy = y + (card_h + gap) * i
        ck = _cat_key(item.get("color"), i)
        _card(slide, x, cy, w, card_h)
        _bar_left(slide, x, cy, card_h, CAT[ck])
        txT = slide.shapes.add_textbox(Inches(x + 0.18), Inches(cy + 0.09),
                                       Inches(w - 0.35), Inches(0.26))
        set_text(txT.text_frame.paragraphs[0], item.get("title", ""),
                 Pt(10), INK, weight="semibold")
        txB = slide.shapes.add_textbox(Inches(x + 0.18), Inches(cy + 0.37),
                                       Inches(w - 0.35), Inches(card_h - 0.46))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tf.paragraphs[0], item.get("body", ""), base_size=Pt(9),
                      base_color=SUB, bold_color=INK, line_spacing=1.25)
    return y + card_h * n + gap * (n - 1)


# =============================================================================
# Section 3: スライドタイプ
# =============================================================================

# 3.1 タイトルスライド（表紙）
def add_title_slide(prs, title, subtitle, date, blank, kpis=None, image_path=None):
    """表紙 — DEEP 濃紺地 + ACCENT チップ + 白 bold タイトル + 短下線 + SKY サブタイトル。

    kpis（任意）: [(値, ラベル, キャプション, 色key), ...]（最大3件推奨）。
        指定時は下部に白カードの KPI カードを並べる（濃紺地では枠線なし）。
    image_path（任意）: 右側に白カードで代表チャート画像を配置する。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, DEEP)

    # ACCENT チップ「APOLLO REPORT」
    chip = _rect(slide, 0.5, 0.52, 1.62, 0.34, ACCENT, round_=True, radius_in=0.04)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = ctf.margin_right = Inches(0.02)
    ctf.margin_top = ctf.margin_bottom = Emu(0)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    crun = cp.add_run()
    crun.text = "APOLLO REPORT"
    crun.font.size = Pt(9.5)
    crun.font.color.rgb = WHITE
    _apply_font(crun, "medium")
    _set_spacing(crun, 140)

    has_right = bool(image_path)
    title_w = 6.6 if has_right else 11.5

    # タイトル（白 bold）
    t_size = Pt(40) if len(title.replace("\n", "")) <= 22 else Pt(32)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(title_w), Inches(1.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], title, t_size, WHITE, line_spacing=1.08, weight="bold")

    # 短下線
    _rect(slide, 0.5, 3.24, 0.62, Emu(38100), ACCENT)

    # サブタイトル（SKY）
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.44), Inches(title_w + 0.1), Inches(0.55))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf2.paragraphs[0], subtitle, Pt(14), SKY, line_spacing=1.25, weight="semibold")

    # 日付（SUB_DK リード位置）
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.20), Inches(6.0), Inches(0.35))
    set_text(txBox3.text_frame.paragraphs[0], date, Pt(11), SUB_DK, weight="light")

    # KPI カード（任意・濃紺地では枠線なしの白カード + 左カテゴリバー）
    if kpis:
        kw, kh, ky = 2.0, 1.18, 5.42
        for i, item in enumerate(kpis[:5]):
            v, l, cap, ck_raw = (list(item) + ["", "", "", None])[:4]
            ck = _cat_key(ck_raw, i)
            x = 0.5 + (kw + 0.18) * i
            _rect(slide, x, ky, kw, kh, CARD, round_=True)
            _bar_left(slide, x, ky, kh, CAT[ck])
            # 値は折返し禁止。長い値はフォントを段階的に縮小して1行に収める
            v_str = str(v)
            v_size = 21 if len(v_str) <= 6 else (17 if len(v_str) <= 8 else 14)
            txV = slide.shapes.add_textbox(Inches(x + 0.16), Inches(ky + 0.12),
                                           Inches(kw - 0.3), Inches(0.45))
            txV.text_frame.word_wrap = False
            set_text(txV.text_frame.paragraphs[0], v_str, Pt(v_size), INK, weight="bold")
            txL = slide.shapes.add_textbox(Inches(x + 0.16), Inches(ky + 0.60),
                                           Inches(kw - 0.3), Inches(0.25))
            set_text(txL.text_frame.paragraphs[0], str(l), Pt(10), SUB, weight="medium")
            if cap:
                txC = slide.shapes.add_textbox(Inches(x + 0.16), Inches(ky + 0.85),
                                               Inches(kw - 0.3), Inches(0.24))
                set_text(txC.text_frame.paragraphs[0], str(cap), Pt(8.5), MUTED, weight="light")

    # 右側チャートカード（任意）
    if has_right:
        mx2, my2, mw2, mh2 = 7.35, 1.05, 5.42, 5.6
        _rect(slide, mx2, my2, mw2, mh2, CARD, round_=True)
        full_path = os.path.join(SNAP, image_path) if not os.path.isabs(image_path) else image_path
        fit_image(slide, full_path, max_x=mx2 + 0.12, max_y=my2 + 0.12,
                  max_w=mw2 - 0.24, max_h=mh2 - 0.24)

    # 下端 ACCENT バー
    _rect(slide, 0, 7.34, 13.333, 0.16, ACCENT)
    return slide


# 3.2 セクション区切り（ライト章扉・ゴースト番号付き）
def add_section_slide(prs, section_num, title, blank, subtitle=None):
    """セクション区切り — BG ライト地 + GHOST 2桁ゴースト数字（右・230pt・Black）。

    ⚠️ 品質ゲートの章扉判定条件を維持していること:
      「^[0-9]{2}$ に一致する2桁数字の単独テキストを持つ／PICTURE・表なし／
       AUTO_SHAPE 4個以下／非空白130字未満」
      本実装: 数字はテキストボックス（AUTO_SHAPE ではない）、AUTO_SHAPE は短下線1個のみ。
    ゴースト数字のみ Noto Sans JP Black を使用する（他は bold/medium）。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, BG)

    try:
        num_text = f"{int(section_num):02d}"
    except (TypeError, ValueError):
        num_text = str(section_num)

    # ゴースト2桁数字（右寄せ・230pt・Black・wrap無効）
    ghost = slide.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.78), Inches(3.9))
    tf_g = ghost.text_frame
    tf_g.word_wrap = False
    tf_g.auto_size = MSO_AUTO_SIZE.NONE
    p_g = tf_g.paragraphs[0]
    p_g.alignment = PP_ALIGN.RIGHT
    run_g = p_g.add_run()
    run_g.text = num_text
    run_g.font.size = Pt(230)
    run_g.font.color.rgb = GHOST
    _apply_font(run_g, "black")

    # "SECTION NN" アイブロウ
    eb = slide.shapes.add_textbox(Inches(0.5), Inches(2.62), Inches(4.0), Inches(0.3))
    ep = eb.text_frame.paragraphs[0]
    set_text(ep, f"SECTION {num_text}", Pt(10), MUTED, weight="medium")
    for _r in ep.runs:
        _set_spacing(_r, 160)

    # 短下線（唯一の AUTO_SHAPE）
    _rect(slide, 0.5, 3.02, 0.62, Emu(38100), ACCENT)

    # セクションタイトル（INK・bold）
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.22), Inches(8.5), Inches(0.95))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], title, Pt(30), INK, line_spacing=1.15, weight="bold")

    # サブ行（省略可）
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.20), Inches(7.6), Inches(0.6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.auto_size = MSO_AUTO_SIZE.NONE
        set_text(tf2.paragraphs[0], subtitle, Pt(13), SUB, weight="medium")
    return slide


# 3.3 チャート+テキスト注釈スライド（主力 — 50%以上）
def add_chart_text_slide(prs, title, sub_message, image_path, annotations, blank,
                         caption=None, chart_label=None, text_side="right",
                         chart_ratio=0.60, source=None, page_num=None, eyebrow=None):
    """チャート主体 + テキスト注釈 — 主力スライドタイプ（§0.9-A0 の主張骨格を載せる）

    Args:
        title: 主張見出し（結論性のある名詞句）
        sub_message: リード文（核心主張の完結した一文・数値込み）
        annotations: 根拠の完結文リスト（各1-2行・最大5項目。最後の1項目は締め文）。
                     文字列リスト → 従来互換の■箇条書き一覧表示。
                     dict リスト（[{"title","body","color"},...]）→ 注釈カードレール
                     （add_annotation_cards）で白カード表示。
        eyebrow: タイトル直上のアイブロウ（章/モジュール名。例 "NEBULA / 環境分析"）
        text_side: "right" or "left"
        chart_ratio: チャート側の幅比率（0.55-0.65）
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title, eyebrow=eyebrow)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    content_w = 12.3
    content_x = 0.5
    gap = 0.3
    chart_w = content_w * chart_ratio - gap / 2
    text_w = content_w * (1 - chart_ratio) - gap / 2
    remaining_h = 6.5 - content_y  # フッターまで使い切る

    if text_side == "right":
        chart_x = content_x
        text_x = content_x + chart_w + gap
    else:
        text_x = content_x
        chart_x = content_x + text_w + gap

    # チャート小見出し
    if chart_label:
        add_chart_label(slide, chart_label, chart_x, content_y, chart_w)
        img_y = content_y + 0.35
        img_h = remaining_h - 0.65
    else:
        img_y = content_y
        img_h = remaining_h - 0.3

    # チャート画像（白カードの上に配置）
    full_path = os.path.join(SNAP, image_path) if not os.path.isabs(image_path) else image_path
    _card_image(slide, full_path, chart_x, img_y, chart_w, img_h)

    # キャプション
    if caption:
        txBox = slide.shapes.add_textbox(Inches(chart_x), Inches(content_y + remaining_h - 0.22),
                                         Inches(chart_w), Inches(0.25))
        set_text(txBox.text_frame.paragraphs[0], caption, Pt(9), MUTED, weight="light")
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # テキスト注釈（dict リストなら注釈カード、文字列リストなら従来の■一覧）
    ann = list(annotations[:5]) if annotations else []
    if ann and isinstance(ann[0], dict):
        add_annotation_cards(slide, text_x, content_y, text_w, ann)
    elif ann:
        add_annotation_block(slide, ann, text_x, content_y, text_w, remaining_h - 0.2)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.4 KPIダッシュボード
def add_kpi_slide(prs, title, sub_message, kpis, blank,
                  source=None, page_num=None):
    """KPIダッシュボード — 白カード + 上部カテゴリ色バー + bold インク数値

    kpis: [{"label":"総特許件数", "value":"1,176", "unit":"件", "trend":"↑12%",
            "color":"blue"(任意カテゴリkey)}, ...]
    4個以下 = 1行配置、5-8個 = 2行配置
    各カード: 上部カテゴリバー + 値(30pt bold インク) + ラベル(SemiBold) + 単位(Light)
    色はカテゴリ順（blue→teal→violet→amber）に固定割当（スライドごとに循環させない）。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(kpis)
    available_w = 11.5
    start_x = 0.9
    gap = 0.2

    if n <= 4:
        n_cols = n
        n_rows = 1
    else:
        n_cols = min(4, (n + 1) // 2)
        n_rows = 2

    card_w = (available_w - gap * (n_cols - 1)) / max(n_cols, 1)
    available_h = 6.5 - content_y
    row_gap = 0.2
    card_h = (available_h - row_gap * (n_rows - 1)) / n_rows
    card_h = min(card_h, 1.9)  # 上限（間延び防止）

    for idx, kpi in enumerate(kpis):
        row = idx // n_cols
        col = idx % n_cols
        x = start_x + col * (card_w + gap)
        y = content_y + row * (card_h + row_gap)
        ck = _cat_key(kpi.get("color"), idx)

        # 白カード + 上部カテゴリ色バー
        _card(slide, x, y, card_w, card_h)
        _bar_top(slide, x, y, card_w, CAT[ck])

        # 値（大数値・bold・インク）+ トレンド
        txV = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.18),
                                       Inches(card_w - 0.4), Inches(0.62))
        p = txV.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = kpi["value"]
        run.font.size = Pt(30)
        run.font.color.rgb = INK
        _apply_font(run, "bold")
        if kpi.get("trend"):
            trend = kpi["trend"]
            if "+" in trend or "↑" in trend or "UP" in trend.upper():
                tc = GOOD_TX
            elif "-" in trend or "↓" in trend or "DOWN" in trend.upper():
                tc = RED_ACCENT
            else:
                tc = MUTED
            run2 = p.add_run()
            run2.text = f" {trend}"
            run2.font.size = Pt(13)
            run2.font.color.rgb = tc
            _apply_font(run2, "medium")

        # ラベル（SemiBold・SUB）
        txL = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.86),
                                       Inches(card_w - 0.4), Inches(0.3))
        set_text(txL.text_frame.paragraphs[0], kpi["label"], Pt(11), SUB, weight="semibold")

        # 単位/キャプション（Light・MUTED・下端）
        if kpi.get("unit"):
            txU = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + card_h - 0.36),
                                           Inches(card_w - 0.4), Inches(0.25))
            set_text(txU.text_frame.paragraphs[0], kpi.get("unit", ""), Pt(9), MUTED, weight="light")

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.4b 結論バンド面（v6 新設）
def add_bands_slide(prs, title, sub_message, bands, blank, page_num=None,
                    conclusion=None, note=None, source=None):
    """結論バンド面 — 淡色ティント地バンド + 左バー + 白抜きチップ + bold 見出し。

    bands: [{"label":"短中期の主軸", "title":"3D NAND ／ SSD", "note":"補足1文",
             "color":"blue"}, ...]（2〜4本）
    conclusion（任意）: 指定時は下部に DEEP 濃紺の結論ピル（白文字・中央寄せ）を置く。
    note（任意）: 指定時は最下部に散文段落（SUB色）を置く。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = max(1, len(bands))
    gap = 0.16
    concl_h = 0.62 if conclusion else 0.0
    note_h = 0.80 if note else 0.0
    bottom_limit = 6.45 - (0.15 if source else 0.0)
    avail = bottom_limit - content_y - (concl_h + 0.12 if conclusion else 0) - (note_h + 0.10 if note else 0)
    bh = min(0.98, max(0.72, (avail - gap * (n - 1)) / n))

    for i, band in enumerate(bands):
        y = content_y + (bh + gap) * i
        ck = _cat_key(band.get("color"), i)
        _rect(slide, 0.5, y, 12.3, bh, TINT[ck], round_=True)
        _bar_left(slide, 0.5, y, bh, CAT[ck])
        if band.get("label"):
            _chip_pill(slide, 0.72, y + bh / 2 - 0.15, band["label"], ck, w=1.5)
        # バンド見出し（bold・インク）
        ty = y + 0.12 if band.get("note") else y
        txT = slide.shapes.add_textbox(Inches(2.45), Inches(ty), Inches(9.6),
                                       Inches(0.42 if band.get("note") else bh))
        tfT = txT.text_frame
        tfT.word_wrap = True
        tfT.auto_size = MSO_AUTO_SIZE.NONE
        if not band.get("note"):
            tfT.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(tfT.paragraphs[0], band.get("title", ""), Pt(16), INK, weight="bold")
        # 注記（SUB）
        if band.get("note"):
            txN = slide.shapes.add_textbox(Inches(2.45), Inches(y + bh - 0.42),
                                           Inches(10.2), Inches(0.34))
            set_text(txN.text_frame.paragraphs[0], band["note"], Pt(9.5), SUB)

    cursor_y = content_y + (bh + gap) * n - gap

    # 結論ピル（DEEP 濃紺・白文字・中央寄せ）
    if conclusion:
        yb = cursor_y + 0.12
        pill = _rect(slide, 1.7, yb, 12.3 - 2.4, concl_h, DEEP, round_=True)
        ptf = pill.text_frame
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.word_wrap = True
        ptf.margin_left = ptf.margin_right = Inches(0.25)
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        set_text(pp, conclusion, Pt(12.5), WHITE, weight="semibold")
        cursor_y = yb + concl_h

    # 下部散文段落
    if note:
        txN = slide.shapes.add_textbox(Inches(0.5), Inches(cursor_y + 0.10),
                                       Inches(12.3), Inches(note_h))
        tfN = txN.text_frame
        tfN.word_wrap = True
        tfN.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tfN.paragraphs[0], note, base_size=Pt(10.5),
                      base_color=SUB, bold_color=INK, line_spacing=1.4)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.5 カードスライド（3-4枚並列）
def add_cards_slide(prs, title, sub_message, cards, blank,
                    source=None, page_num=None):
    """カード並列表示 — 白角丸カード + 上部カテゴリ色バー + インク見出し

    cards: [{"header":"クラスタA", "body":"説明テキスト", "color":"blue"}, ...]
    color はカテゴリキー（"blue"/"teal"/"violet"/"amber"）。旧 API の RGBColor も
    受け付け、近いカテゴリへマップする。省略時はカテゴリ順に固定割当。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(cards)
    gap = 0.25
    total_w = 12.3
    card_w = (total_w - gap * (n - 1)) / max(n, 1)
    card_h = 6.5 - content_y - 0.05  # 下端まで使い切る

    for i, card in enumerate(cards):
        x = 0.5 + i * (card_w + gap)
        ck = _cat_key(card.get("color"), i)

        _card(slide, x, content_y, card_w, card_h)
        _bar_top(slide, x, content_y, card_w, CAT[ck])

        # ヘッダー（インク・SemiBold）
        txH = slide.shapes.add_textbox(Inches(x + 0.18), Inches(content_y + 0.18),
                                       Inches(card_w - 0.36), Inches(0.36))
        set_text(txH.text_frame.paragraphs[0], card["header"], Pt(13.5), INK, weight="semibold")

        # ボディ
        body_y = content_y + 0.60
        txB = slide.shapes.add_textbox(Inches(x + 0.18), Inches(body_y),
                                       Inches(card_w - 0.36), Inches(card_h - 0.78))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        body_text = card.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(4)
                add_rich_runs(p, f"・{item}", base_size=Pt(11), base_color=SUB,
                              bold_color=INK, line_spacing=1.35)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(12),
                          base_color=SUB, bold_color=INK, line_spacing=1.4)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.5b 2×2 マトリクススライド（軸付き・コンサル定番）
def add_matrix_2x2_slide(prs, title, sub_message, x_axis, y_axis, quadrants, blank,
                         source=None, page_num=None):
    """2×2 マトリクス（軸付き）— ティント地象限 + カテゴリ色左バー + インク文字。

    x_axis: {"label": "出願数（活動量）", "low": "少", "high": "多"}
    y_axis: {"label": "成長率（CAGR）", "low": "低", "high": "高"}
    quadrants: 4要素のリスト [左上, 右上, 左下, 右下]
        各要素 {"label": "少数精鋭", "desc": "短い説明", "color": "blue"}（color 省略可）
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    # 描画領域（左に Y 軸ラベル、下に X 軸ラベルの余白を確保）
    left, top, right, bottom = 1.7, content_y + 0.15, 12.6, 6.00
    aw, ah = right - left, bottom - top
    gap = 0.18
    cw, ch = (aw - gap) / 2, (ah - gap) / 2
    pos = [(left, top), (left + cw + gap, top),
           (left, top + ch + gap), (left + cw + gap, top + ch + gap)]
    default_keys = ["blue", "teal", "violet", "amber"]

    for i, q in enumerate(quadrants[:4]):
        qx, qy = pos[i]
        ck = _cat_key(q.get("color"), i) if q.get("color") is not None else default_keys[i]
        _rect(slide, qx, qy, cw, ch, TINT[ck], round_=True)
        _bar_left(slide, qx, qy, ch, CAT[ck])
        tb = slide.shapes.add_textbox(Inches(qx + 0.25), Inches(qy + 0.13),
                                      Inches(cw - 0.45), Inches(ch - 0.26))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # 箱内テキストを上下中央（スカスカ防止・§0.9-G）
        set_text(tf.paragraphs[0], q.get("label", ""), Pt(15), INK, weight="bold")
        if q.get("desc"):
            p = tf.add_paragraph()
            p.space_before = Pt(3)
            add_rich_runs(p, q["desc"], base_size=Pt(11), base_color=SUB,
                          bold_color=INK, line_spacing=1.3)

    # X 軸（下端・右向き矢印 + ラベル + 低/高） — コネクタではなく矢印オートシェイプ
    xa = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(bottom + 0.10),
                                Inches(aw), Inches(0.16))
    xa.fill.solid()
    xa.fill.fore_color.rgb = MUTED
    xa.line.fill.background()
    _shadow_off(xa)
    xl = slide.shapes.add_textbox(Inches(left), Inches(bottom + 0.30), Inches(aw), Inches(0.3))
    set_text(xl.text_frame.paragraphs[0],
             f"{x_axis.get('low','低')}　←　{x_axis.get('label','')}　→　{x_axis.get('high','高')}",
             Pt(11), SUB, weight="medium")
    xl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Y 軸（左端・上向き矢印 + 回転ラベル）
    ya = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(left - 0.38), Inches(top),
                                Inches(0.16), Inches(ah))
    ya.fill.solid()
    ya.fill.fore_color.rgb = MUTED
    ya.line.fill.background()
    _shadow_off(ya)
    yl = slide.shapes.add_textbox(Inches(left - 1.65), Inches(top + ah / 2 - 0.18),
                                  Inches(2.0), Inches(0.36))
    set_text(yl.text_frame.paragraphs[0],
             f"{y_axis.get('low','低')} ← {y_axis.get('label','')} → {y_axis.get('high','高')}",
             Pt(11), SUB, weight="medium")
    yl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    yl.rotation = -90

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.5c 横向き矢羽根フロー（プロセス/因果フロー）
def add_arrow_flow_slide(prs, title, sub_message, steps, blank,
                         source=None, page_num=None):
    """横向きプロセス/因果フロー。先頭=PENTAGON、以降=CHEVRON。

    steps: [{"title":"探索", "desc":"母集団1,176件を取得"}, ...]（3〜6個）
    矢羽根は白文字を載せるため濃色変種（DEEP + CHIP 4色）で塗る。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(steps)
    left, right = 0.5, 12.83
    W = right - left
    overlap = 0.16                      # 矢羽根の先端を次のノッチへ噛み合わせる
    cw = (W + overlap * (n - 1)) / n
    band_h = 1.35
    band_y = content_y + 0.40
    colors = [DEEP, CHIP["blue"], CHIP["teal"], CHIP["violet"], CHIP["amber"]]
    title_size = Pt(16) if n <= 4 else Pt(13)
    desc_top = band_y + band_h + 0.25
    desc_h = max(0.6, 6.35 - desc_top)

    for i, step in enumerate(steps):
        cx = left + i * (cw - overlap)
        color = colors[i % len(colors)]
        shape_type = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        ch = slide.shapes.add_shape(shape_type, Inches(cx), Inches(band_y),
                                    Inches(cw), Inches(band_h))
        ch.fill.solid()
        ch.fill.fore_color.rgb = color
        ch.line.fill.background()
        _shadow_off(ch)

        # ラベル（白・太字・中央。先端を避けて少し左に寄せた領域）
        tb = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(band_y + 0.15),
                                      Inches(cw - 0.70), Inches(band_h - 0.30))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # 箱内テキストを上下中央（スカスカ防止・§0.9-G）
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p, step.get("title", ""), title_size, WHITE, weight="bold")

        # 直下の説明
        if step.get("desc"):
            db = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(desc_top),
                                          Inches(cw - 0.40), Inches(desc_h))
            dtf = db.text_frame
            dtf.word_wrap = True
            dtf.auto_size = MSO_AUTO_SIZE.NONE
            add_rich_runs(dtf.paragraphs[0], step["desc"], base_size=Pt(12),
                          base_color=SUB, bold_color=INK, line_spacing=1.3)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.5d ドーナツチャート（構成比・BLOCK_ARC）
def add_donut_slide(prs, title, sub_message, segments, blank,
                    center_label=None, source=None, page_num=None):
    """ドーナツ図。BLOCK_ARC を構成比に応じた角度で並べる。

    segments: [{"label":"正極材料", "value":330, "color":CAT["blue"]}, ...]（3〜4推奨）
    center_label: ドーナツ中央に置く大きな数値/語（任意）。
    全弧を共通の正方形バウンディングボックスに重ね、rotation=270 で頂点(12時)始まり
    に揃える。終了角>開始角で時計回りにその扇形を塗る（OOXML: swAng=adj2-adj1）。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    total = sum(s.get("value", 0) for s in segments) or 1
    area_top = content_y + 0.20
    area_bottom = 6.30
    R = min(1.95, (area_bottom - area_top) / 2)
    cy = (area_top + area_bottom) / 2
    cx = 3.30
    inner = 0.55
    colors = [CAT["blue"], CAT["teal"], CAT["violet"], CAT["amber"],
              CHIP["blue"], RED_ACCENT]

    DEG = 60000.0 / 100000.0             # 実角度°→adjustment値（python-pptxの格納仕様）
    ang = 0.0                            # 実角度（度）で累積
    for i, seg in enumerate(segments):
        frac = seg.get("value", 0) / total
        ang2 = ang + frac * 360.0
        color = seg.get("color", colors[i % len(colors)])
        if isinstance(color, str):
            color = CAT.get(color, CAT["blue"])
        arc = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC,
                                     Inches(cx - R), Inches(cy - R),
                                     Inches(2 * R), Inches(2 * R))
        arc.adjustments[0] = ang * DEG
        arc.adjustments[1] = (ang2 if ang2 < 359.999 else 359.999) * DEG
        arc.adjustments[2] = inner
        arc.rotation = 270               # 12時始まり（共通中心なので群として回転）
        arc.fill.solid()
        arc.fill.fore_color.rgb = color
        arc.line.color.rgb = BG
        arc.line.width = Pt(1.5)
        _shadow_off(arc)
        ang = ang2

    if center_label:
        cl = slide.shapes.add_textbox(Inches(cx - 1.1), Inches(cy - 0.45),
                                      Inches(2.2), Inches(0.9))
        ctf = cl.text_frame
        ctf.word_wrap = True
        ctf.auto_size = MSO_AUTO_SIZE.NONE
        pc = ctf.paragraphs[0]
        pc.alignment = PP_ALIGN.CENTER
        set_text(pc, center_label, Pt(22), INK, weight="bold")

    # 凡例（右側・色見本 + ラベル + 構成比%）
    lx = cx + R + 0.9
    lw = 12.83 - lx
    n = len(segments)
    row_h = min(0.62, (area_bottom - area_top) / max(n, 1))
    ly0 = cy - (n * row_h) / 2
    for i, seg in enumerate(segments):
        color = seg.get("color", colors[i % len(colors)])
        if isinstance(color, str):
            color = CAT.get(color, CAT["blue"])
        ry = ly0 + i * row_h
        _rect(slide, lx, ry + 0.04, 0.30, 0.30, color, round_=True, radius_in=0.04)
        pct = round(seg.get("value", 0) / total * 100)
        tb = slide.shapes.add_textbox(Inches(lx + 0.45), Inches(ry), Inches(lw - 0.45), Inches(row_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # 箱内テキストを上下中央（スカスカ防止・§0.9-G）
        add_rich_runs(tf.paragraphs[0],
                      f"**{seg.get('label','')}** — {pct}%（{seg.get('value',0)}）",
                      base_size=Pt(13), base_color=SUB, bold_color=INK, line_spacing=1.2)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.5e Issue Tree / ロジックツリー（論証骨格）
def add_issue_tree_slide(prs, title, sub_message, root, branches, blank,
                         source=None, page_num=None):
    """左に論点（根）、右に分解した枝を縦に並べる2階層ロジックツリー。

    root: {"title":"なぜ権利化率に差?", "desc":"上位出願人で2倍超の開き"}
    branches: [{"title":"出願戦略の差", "desc":"量産型 vs 厳選型"}, ...]（2〜5）
    根 = DEEP 濃紺（白文字）、枝 = 白カード + カテゴリ色左バー。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    area_top = content_y + 0.20
    area_bottom = 6.30
    mid_y = (area_top + area_bottom) / 2

    left = 0.5
    root_w = 2.9
    root_h = min(1.7, area_bottom - area_top)
    root_y = mid_y - root_h / 2
    spine_x = left + root_w + 0.95
    arrow_w = 0.55
    child_x = spine_x + arrow_w + 0.05
    child_w = 12.83 - child_x
    n = len(branches)
    gap = 0.22
    child_h = min(1.45, (area_bottom - area_top - gap * (n - 1)) / n)
    total_h = n * child_h + (n - 1) * gap
    start_y = mid_y - total_h / 2
    LINE_T = 0.035                       # 枝線の太さ（細い矩形）

    # 根（論点）ボックス — DEEP 濃紺 + 白文字
    rbox = _rect(slide, left, root_y, root_w, root_h, DEEP, round_=True)
    rtb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(root_y + 0.13),
                                   Inches(root_w - 0.30), Inches(root_h - 0.26))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    rtf.auto_size = MSO_AUTO_SIZE.NONE
    rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(rtf.paragraphs[0], root.get("title", ""), Pt(15), WHITE, weight="bold")
    if root.get("desc"):
        add_rich_runs(rtf.add_paragraph(), root["desc"], base_size=Pt(11),
                      base_color=SUB_DK, bold_color=WHITE, line_spacing=1.25)

    # 根→幹の水平枝（細い矩形）
    _rect(slide, left + root_w, mid_y - LINE_T / 2,
          spine_x - (left + root_w), LINE_T, MUTED)

    child_centers = [start_y + i * (child_h + gap) + child_h / 2 for i in range(n)]
    # 縦の幹（最初〜最後の枝中心を結ぶ細い矩形）
    if n > 1:
        sp_top = min(child_centers[0], mid_y)
        sp_bot = max(child_centers[-1], mid_y)
        _rect(slide, spine_x - LINE_T / 2, sp_top, LINE_T, sp_bot - sp_top, MUTED)

    for i, br in enumerate(branches):
        cyc = child_centers[i]
        cy_box = cyc - child_h / 2
        ck = _cat_key(br.get("color"), i)
        # 幹→枝の右向き矢印（方向表現・オートシェイプ）
        ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(spine_x),
                                    Inches(cyc - 0.11), Inches(arrow_w), Inches(0.22))
        ar.fill.solid()
        ar.fill.fore_color.rgb = CAT[ck]
        ar.line.fill.background()
        _shadow_off(ar)
        # 枝カード（白カード + カテゴリ色左バー）
        _card(slide, child_x, cy_box, child_w, child_h)
        _bar_left(slide, child_x, cy_box, child_h, CAT[ck])
        ctb = slide.shapes.add_textbox(Inches(child_x + 0.28), Inches(cy_box + 0.10),
                                       Inches(child_w - 0.45), Inches(child_h - 0.20))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.auto_size = MSO_AUTO_SIZE.NONE
        set_text(ctf.paragraphs[0], br.get("title", ""), Pt(14), INK, weight="semibold")
        if br.get("desc"):
            add_rich_runs(ctf.add_paragraph(), br["desc"], base_size=Pt(12),
                          base_color=SUB, bold_color=INK, line_spacing=1.25)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.6 プロセスフロー（縦STEP型）
def add_process_slide(prs, title, sub_message, steps, blank,
                      source=None, page_num=None):
    """縦STEPプロセスフロー — 濃色ヘッダー + 白カードボディ

    steps: [{"title":"データ収集", "desc":"特許DBから1,176件を取得"}, ...]
    2個以下 = 大ボックス、3個 = 中、4個以上 = コンパクト
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(steps)
    available_h = 6.5 - content_y
    gap = 0.15
    step_h = (available_h - gap * (n - 1) - 0.5) / n  # 矢印分0.5確保
    step_h = min(step_h, 1.5)

    # フォントサイズ調整
    if n <= 2:
        title_size, desc_size = Pt(16), Pt(14)
    elif n <= 3:
        title_size, desc_size = Pt(14), Pt(13)
    else:
        title_size, desc_size = Pt(13), Pt(12)

    header_w = 2.2
    body_w = 9.8
    colors = [DEEP, CHIP["blue"], CHIP["teal"], CHIP["violet"], CHIP["amber"]]

    for i, step in enumerate(steps):
        sy = content_y + i * (step_h + gap + 0.15)
        color = colors[i % len(colors)]

        # 左ヘッダー（濃色 + 白文字・角丸）
        hdr = _rect(slide, 0.5, sy, header_w, step_h, color, round_=True, radius_in=0.05)
        txH = slide.shapes.add_textbox(Inches(0.6), Inches(sy + 0.08),
                                       Inches(header_w - 0.2), Inches(step_h - 0.16))
        tf_h = txH.text_frame
        tf_h.word_wrap = True
        tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_h = tf_h.paragraphs[0]
        p_h.alignment = PP_ALIGN.CENTER
        set_text(p_h, f"STEP {i+1}", Pt(10), SUB_DK if color == DEEP else WHITE, weight="medium")
        p_t = tf_h.add_paragraph()
        p_t.alignment = PP_ALIGN.CENTER
        set_text(p_t, step["title"], title_size, WHITE, weight="bold")

        # 右ボディ（白カード）
        _card(slide, 0.5 + header_w + 0.1, sy, body_w, step_h)
        txB = slide.shapes.add_textbox(Inches(0.5 + header_w + 0.28), Inches(sy + 0.10),
                                       Inches(body_w - 0.36), Inches(step_h - 0.2))
        tf_b = txB.text_frame
        tf_b.word_wrap = True
        tf_b.auto_size = MSO_AUTO_SIZE.NONE
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_rich_runs(tf_b.paragraphs[0], step.get("desc", ""),
                      base_size=desc_size, base_color=SUB, bold_color=INK, line_spacing=1.3)

        # 下矢印（最後のステップ以外）
        if i < n - 1:
            arrow_y = sy + step_h + 0.02
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, Inches(1.2), Inches(arrow_y),
                Inches(0.5), Inches(gap + 0.05)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()
            _shadow_off(arrow)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.7 ステップアップスライド（階段型ロードマップ）
def add_stepup_slide(prs, title, sub_message, phases, blank,
                     source=None, page_num=None):
    """ステップアップ（階段型ロードマップ）— 濃色ヘッダー + 白カードボディ

    phases: [{"header":"短期", "body":"基盤構築", "color":"blue"}, ...]
    左から右へ棒の高さが上がる。3-4段を推奨。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(phases)
    gap = 0.2
    total_w = 12.3
    bar_w = (total_w - gap * (n - 1)) / max(n, 1)
    base_y = 6.5  # フッター直上
    max_h = base_y - content_y - 0.2
    strong = [CHIP["blue"], CHIP["teal"], CHIP["violet"], CHIP["amber"]]

    for i, phase in enumerate(phases):
        x = 0.5 + i * (bar_w + gap)
        # 高さを段階的に上げる（最小50%、最大100%）
        ratio = 0.5 + 0.5 * (i / max(n - 1, 1))
        bar_h = max_h * ratio
        y = base_y - bar_h
        color = phase.get("color")
        if color is None or isinstance(color, str):
            color = strong[CAT_ORDER.index(_cat_key(color, i))]

        # ヘッダー部（上部、濃色 + 白文字）
        header_h = min(0.5, bar_h * 0.25)
        _rect(slide, x, y, bar_w, header_h, color)
        txH = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05),
                                       Inches(bar_w - 0.2), Inches(header_h - 0.1))
        set_text(txH.text_frame.paragraphs[0], phase["header"], Pt(14), WHITE, weight="bold")
        txH.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # ボディ部（下部、白カード）
        body_y = y + header_h
        body_h = bar_h - header_h
        _rect(slide, x, body_y, bar_w, body_h, CARD, CARD_LINE)
        txB = slide.shapes.add_textbox(Inches(x + 0.15), Inches(body_y + 0.12),
                                       Inches(bar_w - 0.3), Inches(body_h - 0.24))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        body_text = phase.get("body", "")
        if isinstance(body_text, list):
            for j, item in enumerate(body_text):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_rich_runs(p, f"・{item}", base_size=Pt(11), base_color=SUB,
                              bold_color=INK, line_spacing=1.3)
        else:
            add_rich_runs(tf.paragraphs[0], body_text, base_size=Pt(12),
                          base_color=SUB, bold_color=INK, line_spacing=1.3)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.8 比較スライド（左 vs 右）
def add_compare_slide(prs, title, sub_message, left_title, left_items,
                      right_title, right_items, blank,
                      left_color=ACCENT, right_color=RED_ACCENT,
                      source=None, page_num=None):
    """左右比較スライド — 色ヘッダーバー + 白カードボディ

    left_items / right_items: 各3-5項目の短い注釈リスト
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    col_w = 5.7
    left_x = 0.5
    right_x = 6.9
    remaining_h = 6.5 - content_y
    header_h = 0.42

    # 中央 "VS" マーカー
    vs_box = slide.shapes.add_textbox(Inches(6.1), Inches(content_y + 1.5),
                                      Inches(1.0), Inches(0.5))
    p_vs = vs_box.text_frame.paragraphs[0]
    p_vs.alignment = PP_ALIGN.CENTER
    set_text(p_vs, "VS", Pt(16), MUTED, weight="semibold")

    # 中央区切り線（ヘアライン）
    _rect(slide, 6.55, content_y, HAIRLINE, remaining_h, CARD_LINE)

    for side_x, side_title, side_items, side_color in [
        (left_x, left_title, left_items, left_color),
        (right_x, right_title, right_items, right_color),
    ]:
        # カラムヘッダー（色バー + 白文字）
        _rect(slide, side_x, content_y, col_w, header_h, side_color,
              round_=True, radius_in=0.05)
        txBox = slide.shapes.add_textbox(Inches(side_x + 0.1), Inches(content_y + 0.04),
                                         Inches(col_w - 0.2), Inches(header_h - 0.08))
        set_text(txBox.text_frame.paragraphs[0], side_title, Pt(15), WHITE, weight="bold")
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # ボディ（白カード + 注釈）
        body_y = content_y + header_h + 0.10
        body_h = remaining_h - header_h - 0.25
        _card(slide, side_x, body_y, col_w, body_h)
        add_annotation_block(slide, side_items, side_x + 0.12, body_y + 0.10,
                             col_w - 0.24, body_h - 0.2, font_size=13)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.9 テーブルスライド
def add_table_slide(prs, title, sub_message, headers, rows, blank,
                    col_widths=None, highlight_rows=None, annotations=None,
                    source=None, page_num=None, cell_colors=None):
    """テーブル + オプション注釈テキスト — DEEP ヘッダー + ゼブラ + 罫線レス

    highlight_rows: ハイライト行のインデックスリスト（青ティント地）
    annotations: テーブル横に注釈テキスト表示
    cell_colors: {(row_idx, col_idx): 'good'|'part'} — 判定セルの文字色指定。
        row_idx はデータ行の0始まりインデックス。'good'=GOOD_TX(緑)、
        'part'=PART_TX(琥珀)。指定セルは SemiBold で描画される。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n_cols = len(headers)
    n_rows = len(rows) + 1

    if annotations:
        table_w = 7.5
        text_x = 8.3
        text_w = 4.5
    else:
        table_w = 12.3
        text_x = None
        text_w = 0

    # 行高を残余スペースに合わせて動的計算
    available_table_h = 6.4 - content_y
    row_h = min(0.55, max(0.35, available_table_h / n_rows))
    table_h = row_h * n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(content_y), Inches(table_w), Inches(table_h)
    )
    table = table_shape.table
    _plain_table_style(table)
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # DEEP ヘッダー行
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(cell.text_frame.paragraphs[0], header, Pt(12.5), WHITE, weight="semibold")
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP

    # ゼブラ（CARD/BG 交互）データ行
    highlight_rows = highlight_rows or []
    cell_colors = cell_colors or {}
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            judge = cell_colors.get((i, j))
            if judge == 'good':
                set_text(cell.text_frame.paragraphs[0], str(value), Pt(12),
                         GOOD_TX, weight="semibold")
            elif judge == 'part':
                set_text(cell.text_frame.paragraphs[0], str(value), Pt(12),
                         PART_TX, weight="semibold")
            else:
                set_text(cell.text_frame.paragraphs[0], str(value), Pt(12), INK)
            cell.fill.solid()
            if i in highlight_rows:
                cell.fill.fore_color.rgb = TINT["blue"]
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = CARD
            else:
                cell.fill.fore_color.rgb = BG

    _strip_table_borders(table)

    # 注釈テキスト（テーブル横・白カード）
    if annotations and text_x:
        remaining_h = 6.4 - content_y
        _card(slide, text_x, content_y, text_w, remaining_h)
        add_annotation_block(slide, annotations, text_x + 0.08, content_y + 0.08,
                             text_w - 0.16, remaining_h - 0.16, font_size=12)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.10 プログレスバースライド
def add_progress_bar_slide(prs, title, sub_message, items, blank,
                           source=None, page_num=None):
    """水平プログレスバー

    items: [{"label":"クラスタA", "value":58, "max_value":100, "color":"blue"}, ...]
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(items)
    available_h = 6.5 - content_y
    bar_gap = 0.1
    bar_group_h = (available_h - 0.2) / n
    bar_h = min(0.42, bar_group_h * 0.5)
    label_h = bar_group_h - bar_h - bar_gap
    bar_max_w = 9.0

    for i, item in enumerate(items):
        gy = content_y + i * bar_group_h
        ck = _cat_key(item.get("color"), i)
        pct = item["value"] / max(item.get("max_value", 100), 1)
        bar_w = bar_max_w * pct

        # ラベル（左・SemiBold インク・バーと同じ行に上下中央で配置）
        txL = slide.shapes.add_textbox(Inches(0.5), Inches(gy + label_h - 0.06),
                                       Inches(2.9), Inches(bar_h + 0.12))
        txL.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        txL.text_frame.word_wrap = True
        set_text(txL.text_frame.paragraphs[0], item["label"], Pt(13), INK, weight="semibold")

        # 背景バー（GHOST・全幅）
        _rect(slide, 3.5, gy + label_h, bar_max_w, bar_h, GHOST, round_=True, radius_in=0.05)

        # 値バー（カテゴリ色）
        if bar_w > 0.1:
            _rect(slide, 3.5, gy + label_h, bar_w, bar_h, CAT[ck],
                  round_=True, radius_in=0.05)

        # パーセンテージ（バーの右端・CHIP 濃色文字）
        txP = slide.shapes.add_textbox(Inches(3.5 + bar_w + 0.1), Inches(gy + label_h),
                                       Inches(1.5), Inches(bar_h))
        txP.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(txP.text_frame.paragraphs[0],
                 f"{item['value']}{item.get('unit', '%')}", Pt(13), CHIP[ck], weight="semibold")

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.11 トライアングルスライド（3要素の関係図）
def add_triangle_slide(prs, title, sub_message, elements, blank,
                       source=None, page_num=None):
    """3要素トライアングル関係図 — 濃色ヘッダー + 白カードボディ

    elements: [
        {"title":"技術", "body":"SiC/GaN半導体", "color":"blue"},
        {"title":"市場", "body":"EV・再エネ需要", "color":"teal"},
        {"title":"政策", "body":"グリーン成長戦略", "color":"violet"},
    ]
    上1 + 下2 の三角配置 + 関係線
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    card_w = 3.5
    card_h = 2.0

    # 三角の3頂点座標
    positions = [
        (5.0, content_y + 0.2),          # 上中央
        (1.5, content_y + 2.8),          # 左下
        (8.5, content_y + 2.8),          # 右下
    ]

    # 関係線（3辺・カードの下に描くため先に描画）
    arrow_pairs = [(0, 1), (0, 2), (1, 2)]
    for a, b in arrow_pairs:
        ax = positions[a][0] + card_w / 2
        ay = positions[a][1] + card_h
        bx = positions[b][0] + card_w / 2
        by = positions[b][1]
        if b == 2 and a == 1:
            ay = positions[a][1] + card_h / 2
            by = positions[b][1] + card_h / 2
        # コネクタ（p:cxnSp）はファイル破損を起こすため、回転矩形で関係線を描く
        _add_line(slide, ax, ay, bx, by, MUTED, weight_pt=1.0)

    for i, (elem, (px, py)) in enumerate(zip(elements[:3], positions)):
        ck = _cat_key(elem.get("color"), i)

        # カードヘッダー（濃色 + 白文字）
        _rect(slide, px, py, card_w, 0.45, CHIP[ck], round_=True, radius_in=0.05)
        txH = slide.shapes.add_textbox(Inches(px + 0.1), Inches(py + 0.05),
                                       Inches(card_w - 0.2), Inches(0.35))
        set_text(txH.text_frame.paragraphs[0], elem["title"], Pt(14), WHITE, weight="bold")
        txH.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # カードボディ（白カード）
        _card(slide, px, py + 0.50, card_w, card_h - 0.50)
        txB = slide.shapes.add_textbox(Inches(px + 0.18), Inches(py + 0.62),
                                       Inches(card_w - 0.36), Inches(card_h - 0.76))
        tf = txB.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        add_rich_runs(tf.paragraphs[0], elem.get("body", ""),
                      base_size=Pt(12), base_color=SUB, bold_color=INK, line_spacing=1.4)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.12 ピラミッドスライド
def add_pyramid_slide(prs, title, sub_message, levels, blank,
                      source=None, page_num=None):
    """ピラミッド（上が小、下が大の台形積み重ね）— 濃色変種 + 白文字

    levels: [{"title":"萌芽技術", "detail":"ノイズ6テーマ"}, ...]  上→下の順
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(levels)
    total_h = 6.5 - content_y - 0.2
    level_h = total_h / n
    base_w = 10.0
    center_x = 6.66  # スライド中央
    colors = [CHIP["amber"], CHIP["teal"], CHIP["blue"], CHIP["violet"], DEEP,
              CHIP["teal"]]

    for i, level in enumerate(levels):
        ratio_top = (i + 0.3) / n
        ratio_bot = (i + 1.3) / n
        lw = base_w * (ratio_top + ratio_bot) / 2
        lx = center_x - lw / 2
        ly = content_y + i * level_h
        color = colors[i % len(colors)]

        trap = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID, Inches(lx), Inches(ly),
            Inches(lw), Inches(level_h - 0.05)
        )
        trap.fill.solid()
        trap.fill.fore_color.rgb = color
        trap.line.fill.background()
        _shadow_off(trap)

        txBox = slide.shapes.add_textbox(
            Inches(lx + 0.3), Inches(ly + 0.1),
            Inches(lw - 0.6), Inches(level_h - 0.2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p, level["title"], Pt(14), WHITE, weight="bold")
        if level.get("detail"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            set_text(p2, level["detail"], Pt(11), WHITE)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.13 仮説検証スライド
def add_hypothesis_slide(prs, title, sub_message, hypotheses, blank,
                         source=None, page_num=None):
    """仮説検証テーブル — DEEP ヘッダー + ゼブラ + 判定は色文字（背景に色を載せない）

    hypotheses: [
        {"id":"H1", "hypothesis":"A社は3年以内にシェア首位", "verdict":"partially",
         "evidence":"シェア2位に浮上も首位とのギャップは依然5%"},
        ...
    ]
    verdict: "confirmed" -> ✓ 支持 (緑), "rejected" -> ✕ 不支持 (赤),
             "partially" -> △ 部分支持 (琥珀)
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    headers = ["ID", "仮説", "判定", "エビデンス"]
    n_rows = len(hypotheses) + 1
    available_h = 6.4 - content_y
    row_h = min(0.55, max(0.40, available_h / n_rows))
    table_h = row_h * n_rows

    VERDICT_MAP = {
        "confirmed": ("✓ 支持", GOOD_TX),
        "rejected": ("✕ 不支持", RED_ACCENT),
        "partially": ("△ 部分支持", PART_TX),
    }

    table_shape = slide.shapes.add_table(
        n_rows, 4, Inches(0.5), Inches(content_y), Inches(12.3), Inches(table_h)
    )
    table = table_shape.table
    _plain_table_style(table)
    # 列幅: ID=0.8, 仮説=4.3, 判定=1.6, エビデンス=5.6
    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(4.3)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(5.6)

    # DEEP ヘッダー
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(cell.text_frame.paragraphs[0], header, Pt(12.5), WHITE, weight="semibold")
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP

    # データ行（ゼブラ CARD/BG）
    for i, hyp in enumerate(hypotheses):
        verdict_key = hyp.get("verdict", "partially")
        verdict_label, verdict_color = VERDICT_MAP.get(verdict_key, ("△ 部分支持", PART_TX))

        row_data = [hyp.get("id", ""), hyp.get("hypothesis", ""),
                    verdict_label, hyp.get("evidence", "")]

        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if j == 2:
                # 判定セルは色文字（SemiBold）— 背景に色は載せない
                set_text(cell.text_frame.paragraphs[0], val, Pt(11.5),
                         verdict_color, weight="semibold")
            else:
                set_text(cell.text_frame.paragraphs[0], str(val), Pt(12), INK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i % 2 == 0 else BG

    _strip_table_borders(table)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.14 タイムラインスライド
def add_timeline_slide(prs, title, sub_message, events, blank,
                       source=None, page_num=None):
    """水平タイムライン — 白カードのイベントラベル + カテゴリ色マーカー

    events: [{"year":"2015", "title":"CNF政策支援開始", "color":"blue"}, ...]
    ラベルは上下交互配置。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(events)
    total_w = 11.5
    start_x = 0.9
    line_y = content_y + 1.9  # タイムライン中心位置

    # 水平線（インク）
    _rect(slide, start_x, line_y, total_w, Emu(19050), INK)

    step = total_w / max(n - 1, 1) if n > 1 else total_w
    for i, ev in enumerate(events):
        x = start_x + i * step
        ck = _cat_key(ev.get("color"), i)
        color = CAT[ck]

        # マーカー円
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.12), Inches(line_y - 0.12),
            Inches(0.30), Inches(0.30)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = CARD
        dot.line.width = Emu(19050)
        _shadow_off(dot)

        # 年ラベル（マーカー上、インク太字）
        txY = slide.shapes.add_textbox(Inches(x - 0.37), Inches(line_y - 0.46),
                                       Inches(0.8), Inches(0.25))
        p = txY.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p, ev["year"], Pt(10), INK, weight="semibold")

        # イベントラベル（白カード・交互に上下配置して重なり回避）
        card_w, card_h = 1.7, 0.78
        if i % 2 == 0:
            ty = line_y + 0.45
        else:
            ty = line_y - 0.45 - card_h - 0.25
        cx0 = x - card_w / 2 + 0.03
        _card(slide, cx0, ty, card_w, card_h)
        _bar_left(slide, cx0, ty, card_h, color)
        txE = slide.shapes.add_textbox(Inches(cx0 + 0.12), Inches(ty + 0.06),
                                       Inches(card_w - 0.2), Inches(card_h - 0.12))
        tf = txE.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        set_text(p2, ev["title"], Pt(9), SUB, line_spacing=1.15)

        # 縦線（マーカーからカードへ届くまで引く）
        if i % 2 == 0:
            vline_y = line_y + 0.18
            vline_h = 0.27          # マーカー下端 → 下側カード上端
        else:
            vline_y = ty + card_h
            vline_h = (line_y - 0.12) - (ty + card_h)  # 上側カード下端 → マーカー上端
        if vline_h > 0:
            _rect(slide, x + 0.02, vline_y, Emu(9525), vline_h, color)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 3.15 クロージングスライド
def add_closing_slide(prs, report_title, blank, date_str=None):
    """クロージング（裏表紙）— DEEP 濃紺地。"Thank You" は置かない。

    短下線 + レポートタイトル(bold) + 「APOLLO Patent Analytics Platform」+
    日付（任意）+ 下端 ACCENT バー。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, DEEP)

    # 短下線（中央）
    _rect(slide, 5.92, 2.62, 1.5, Emu(38100), ACCENT)

    # レポートタイトル（白 bold・中央寄せ）
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.33), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    set_text(tf.paragraphs[0], report_title, Pt(28), WHITE, weight="bold")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # APOLLO ブランディング
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.90), Inches(10.33), Inches(0.4))
    p2 = txBox2.text_frame.paragraphs[0]
    set_text(p2, "APOLLO Patent Analytics Platform", Pt(11), SUB_DK, weight="medium")
    p2.alignment = PP_ALIGN.CENTER
    for _r in p2.runs:
        _set_spacing(_r, 100)

    # 日付（任意）
    if date_str:
        txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.33), Inches(10.33), Inches(0.35))
        p3 = txBox3.text_frame.paragraphs[0]
        set_text(p3, str(date_str), Pt(10), MUT_DK, weight="light")
        p3.alignment = PP_ALIGN.CENTER

    # 下端 ACCENT バー
    _rect(slide, 0, 7.34, 13.333, 0.16, ACCENT)
    return slide


# =============================================================================
# 補助スライドタイプ
# =============================================================================

# 目次スライド
def add_toc_slide(prs, title, items, blank, page_num=None):
    """目次スライド — 白カード行の目次

    items = [{"num":1, "title":"セクション名", "page":"P5"}, ...]
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)

    n = len(items)
    table_x, table_y, table_w = 1.5, sub_y + 0.15, 10.3
    row_h = min(0.62, (6.5 - table_y) / max(n, 1))
    card_h = row_h - 0.10

    for i, item in enumerate(items):
        y = table_y + i * row_h
        _card(slide, table_x, y, table_w, card_h)

        # 番号（ACCENT）
        txNum = slide.shapes.add_textbox(Inches(table_x + 0.25), Inches(y),
                                         Inches(0.8), Inches(card_h))
        txNum.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(txNum.text_frame.paragraphs[0],
                 f"{item.get('num', i+1):02d}", Pt(13), ACCENT, weight="semibold")

        # セクション名（インク SemiBold）
        txTitle = slide.shapes.add_textbox(Inches(table_x + 1.2), Inches(y),
                                           Inches(7.0), Inches(card_h))
        txTitle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(txTitle.text_frame.paragraphs[0], item["title"], Pt(13), INK, weight="semibold")

        # ページ番号（MUTED）
        txPage = slide.shapes.add_textbox(Inches(table_x + 8.5), Inches(y),
                                          Inches(1.5), Inches(card_h))
        txPage.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = txPage.text_frame.paragraphs[0]
        set_text(p, item.get("page", ""), Pt(12), MUTED, weight="light")
        p.alignment = PP_ALIGN.RIGHT

    add_bottom_bar_and_footer(slide, page_num)
    return slide


# デュアルパネルスライド（2カラムチャート比較）
def add_dual_panel_slide(prs, title, sub_message,
                         left_label, left_image, left_caption,
                         right_label, right_image, right_caption,
                         left_bullets=None, right_bullets=None,
                         blank=None, source=None, page_num=None):
    """2カラムチャート — 2つの可視化を白カードで並列比較"""
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    col_w = 5.9
    left_x = 0.5
    right_x = 6.9
    remaining_h = 6.5 - content_y

    if left_bullets or right_bullets:
        n_bullets = max(len(left_bullets or []), len(right_bullets or []))
        if n_bullets <= 1:
            chart_h = remaining_h * 0.78
        elif n_bullets <= 2:
            chart_h = remaining_h * 0.68
        else:
            chart_h = remaining_h * 0.58
        text_y = content_y + chart_h + 0.1
        text_h = remaining_h - chart_h - 0.3
    else:
        chart_h = remaining_h - 0.5
        text_y = None
        text_h = 0

    for side_x, label, img_path, caption, bullets in [
        (left_x, left_label, left_image, left_caption, left_bullets),
        (right_x, right_label, right_image, right_caption, right_bullets),
    ]:
        add_chart_label(slide, label, side_x, content_y, col_w)
        full_path = os.path.join(SNAP, img_path) if not os.path.isabs(img_path) else img_path
        _card_image(slide, full_path, side_x, content_y + 0.35, col_w, chart_h - 0.40)

        if caption:
            txBox = slide.shapes.add_textbox(Inches(side_x), Inches(content_y + chart_h),
                                             Inches(col_w), Inches(0.25))
            set_text(txBox.text_frame.paragraphs[0], caption, Pt(9), MUTED, weight="light")

        if bullets and text_y:
            add_annotation_block(slide, bullets, side_x, text_y, col_w, text_h, font_size=13)

    # 中央区切り線（ヘアライン）
    _rect(slide, 6.55, content_y, HAIRLINE, remaining_h, CARD_LINE)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# ナラティブスライド（テキスト主体 — サマリー/結論限定）
def add_narrative_slide(prs, title, sub_message, paragraphs, blank,
                        source=None, page_num=None):
    """テキスト主体 — エグゼクティブサマリーと結論にのみ使用

    全スライドの10%以下に制限すること。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y) if sub_message else sub_y + 0.1

    remaining_h = 6.5 - content_y
    _card(slide, 0.5, content_y, 12.3, remaining_h - 0.15)
    add_annotation_block(slide, paragraphs, 0.75, content_y + 0.20,
                         11.8, remaining_h - 0.55, font_size=15)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# チャート全画面スライド
def add_image_slide(prs, title, sub_message, image_path, blank,
                    caption=None, chart_label=None, source=None, page_num=None):
    """チャート全画面 — 画像が主役のスライド（白カードの上に配置）"""
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    if chart_label:
        add_chart_label(slide, chart_label, 0.5, content_y, 12.3)
        img_y = content_y + 0.35
    else:
        img_y = content_y

    full_path = os.path.join(SNAP, image_path) if not os.path.isabs(image_path) else image_path
    img_h = (6.40 if not caption else 6.20) - img_y
    _card_image(slide, full_path, 0.5, img_y, 12.3, img_h)

    if caption:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.28), Inches(12.3), Inches(0.25))
        set_text(txBox.text_frame.paragraphs[0], caption, Pt(9.5), MUTED, weight="light")
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 推奨アクションスライド
def add_recommendation_slide(prs, title, sub_message, recommendations, blank,
                             source=None, page_num=None):
    """推奨アクション — 白カード + 優先度色左バー

    recommendations: [{"priority":"高","title":"出願強化","timeframe":"短期","desc":"..."},...]
    """
    PRIORITY_COLORS = {"高": RED_ACCENT, "中": CAT["amber"], "低": CAT["teal"]}
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    n = len(recommendations)
    available_h = 6.5 - content_y
    card_h = min(1.3, (available_h - 0.1 * (n - 1)) / n)

    for i, rec in enumerate(recommendations):
        y = content_y + i * (card_h + 0.1)
        p_color = PRIORITY_COLORS.get(rec.get("priority", "中"), CAT["amber"])

        _card(slide, 0.5, y, 12.33, card_h)
        _bar_left(slide, 0.5, y, card_h, p_color)

        txBox_p = slide.shapes.add_textbox(Inches(0.78), Inches(y + 0.08), Inches(0.8), Inches(0.3))
        set_text(txBox_p.text_frame.paragraphs[0], f"[{rec.get('priority', '中')}]",
                 Pt(10), p_color, weight="semibold")

        txBox_t = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.06), Inches(5.8), Inches(0.34))
        set_text(txBox_t.text_frame.paragraphs[0], rec["title"], Pt(15), INK, weight="bold")

        if rec.get("timeframe"):
            txBox_tf = slide.shapes.add_textbox(Inches(9.0), Inches(y + 0.10), Inches(3.5), Inches(0.3))
            p_tf = txBox_tf.text_frame.paragraphs[0]
            set_text(p_tf, rec["timeframe"], Pt(12), MUTED, weight="light")
            p_tf.alignment = PP_ALIGN.RIGHT

        if rec.get("desc"):
            txBox_d = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.44),
                                               Inches(10.5), Inches(card_h - 0.54))
            tf_d = txBox_d.text_frame
            tf_d.word_wrap = True
            tf_d.auto_size = MSO_AUTO_SIZE.NONE
            add_rich_runs(tf_d.paragraphs[0], rec["desc"], base_size=Pt(12),
                          base_color=SUB, bold_color=INK, line_spacing=1.3)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# =============================================================================
# 章構成スライド（参考デッキ準拠・主張骨格 §0.9-A0 を体現）
# =============================================================================

# 重心移動スライド（PAST → PRESENT・エグゼクティブサマリー型）
def add_shift_slide(prs, title, lead, past, present, closing, blank,
                    eyebrow=None, source=None, page_num=None):
    """重心移動（PAST → PRESENT）— 過去の主役から現在の重点への移行を1枚で語る。

    左カード「PAST ・ 過去の主役」、中央に右向き矢印、右カード
    「PRESENT ・ 現在の重点」を同サイズで並べ、下部に頑健性の締め文を置く。

    Args:
        title: 主張見出し（結論性のある名詞句）
        lead: リード文（タイトル直下のプレーンなリード文。`add_sub_message` 使用）
        past / present: dict `{"label":"PAST ・ 過去の主役", "heading":"短い名詞句",
                              "desc":"1〜2文"}`
        closing: 締め文（下部に地の文一文。「4つの独立した手法がいずれも同じ
                 方向を指す、頑健な結論」等）
        eyebrow: タイトル直上のアイブロウ（章/モジュール名）
    主張骨格（§0.9-A0）: アイブロウ→主張見出し→リード文→根拠（PAST/PRESENT 2枚）
                        →締め文（So What）。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title, eyebrow=eyebrow)
    content_y = add_sub_message(slide, lead, y=sub_y)

    # 締め文を下部に確保し、その上にカード帯を置く
    closing_h = 0.75 if closing else 0.0
    band_top = content_y + 0.10
    band_bottom = 6.45 - (closing_h + 0.15 if closing else 0.0)
    band_h = band_bottom - band_top

    arrow_w = 1.05
    gap = 0.30
    total_w = 12.3
    card_w = (total_w - arrow_w - gap * 2) / 2
    left_x = 0.5
    arrow_x = left_x + card_w + gap
    right_x = arrow_x + arrow_w + gap

    for cx, data, accent in [(left_x, past, MUTED), (right_x, present, ACCENT)]:
        _card(slide, cx, band_top, card_w, band_h)
        # 上端アクセント帯（PAST=ミュート／PRESENT=ACCENT で現在を際立たせる）
        _bar_top(slide, cx, band_top, card_w, accent)

        tb = slide.shapes.add_textbox(Inches(cx + 0.25), Inches(band_top + 0.18),
                                      Inches(card_w - 0.50), Inches(band_h - 0.36))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # 上下中央（スカスカ防止・§0.9-G）
        # label（小・ミュート）
        set_text(tf.paragraphs[0], data.get("label", ""), Pt(11), MUTED, weight="medium")
        # heading（bold インク）
        ph = tf.add_paragraph()
        ph.space_before = Pt(4)
        set_text(ph, data.get("heading", ""), Pt(20), INK, weight="bold", line_spacing=1.15)
        # desc（Regular）
        if data.get("desc"):
            pd = tf.add_paragraph()
            pd.space_before = Pt(6)
            add_rich_runs(pd, data["desc"], base_size=Pt(13),
                          base_color=SUB, bold_color=INK, line_spacing=1.35)

    # 中央の右向き矢印（インク・オートシェイプ。コネクタ禁止）
    arr_h = 0.70
    arr = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(band_top + band_h / 2 - arr_h / 2),
        Inches(arrow_w), Inches(arr_h)
    )
    arr.fill.solid()
    arr.fill.fore_color.rgb = INK
    arr.line.fill.background()
    _shadow_off(arr)

    # 締め文（下部・青ティント + ACCENT 左バー）
    if closing:
        _rect(slide, 0.5, band_bottom + 0.15, 12.3, closing_h, TINT["blue"], round_=True)
        _bar_left(slide, 0.5, band_bottom + 0.15, closing_h, ACCENT)
        ctb = slide.shapes.add_textbox(Inches(0.78), Inches(band_bottom + 0.15),
                                       Inches(11.9), Inches(closing_h))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.auto_size = MSO_AUTO_SIZE.NONE
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_rich_runs(ctf.paragraphs[0], closing, base_size=Pt(13),
                      base_color=INK, bold_color=INK, line_spacing=1.35, weight="medium")

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# クロス統合スライド（N手法が同じ結論へ収束）
def add_convergence_slide(prs, title, methods, conclusion, blank,
                          eyebrow=None, source=None, page_num=None):
    """クロス統合（N手法 → 1つの頑健な結論へ収束）。

    左に N 個の手法行（手法名→その手法の発見を1文）、各行から細い矢印が
    右の大きな「頑健な結論」ボックスへ収束する。

    Args:
        methods: list of `{"method":"俯瞰図分析", "finding":"…1文"}`（3〜5件）
        conclusion: dict `{"headline":"…結論見出し", "detail":"…手法に依存しない構造変化"}`
        eyebrow: タイトル直上のアイブロウ
    主張骨格（§0.9-A0）: 複数手法を束ねた「束ねスライド」（§0.9-C）。各手法の発見＝
                        根拠、右の結論箱＝主張見出し＋So What。
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title, eyebrow=eyebrow)
    # リード文は結論見出しを流用して主張骨格を満たす
    content_y = add_sub_message(slide, conclusion.get("headline", ""), y=sub_y)

    area_top = content_y + 0.15
    area_bottom = 6.45
    area_h = area_bottom - area_top

    left_x = 0.5
    method_w = 5.6
    arrow_w = 0.55
    concl_x = left_x + method_w + arrow_w + 0.35
    concl_w = 12.83 - concl_x

    n = len(methods)
    gap = 0.18
    row_h = min(1.30, (area_h - gap * (n - 1)) / max(n, 1))
    total_rows_h = row_h * n + gap * (n - 1)
    start_y = area_top + (area_h - total_rows_h) / 2

    row_centers = []
    for i, m in enumerate(methods):
        ry = start_y + i * (row_h + gap)
        row_centers.append(ry + row_h / 2)
        ck = _cat_key(m.get("color"), i)
        # 手法行（白カード + カテゴリ色左バー + 手法名 + 発見文）
        _card(slide, left_x, ry, method_w, row_h)
        _bar_left(slide, left_x, ry, row_h, CAT[ck])
        tb = slide.shapes.add_textbox(Inches(left_x + 0.28), Inches(ry + 0.08),
                                      Inches(method_w - 0.45), Inches(row_h - 0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(tf.paragraphs[0], m.get("method", ""), Pt(13), INK, weight="semibold")
        if m.get("finding"):
            pf = tf.add_paragraph()
            pf.space_before = Pt(2)
            add_rich_runs(pf, m["finding"], base_size=Pt(11),
                          base_color=SUB, bold_color=INK, line_spacing=1.25)

    # 結論ボックス（青ティント地 + ACCENT 上端帯）
    _rect(slide, concl_x, area_top, concl_w, area_h, TINT["blue"], round_=True)
    _rect(slide, concl_x + 0.10, area_top, concl_w - 0.20, Emu(64008), ACCENT)
    ctb = slide.shapes.add_textbox(Inches(concl_x + 0.28), Inches(area_top + 0.25),
                                   Inches(concl_w - 0.56), Inches(area_h - 0.50))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.auto_size = MSO_AUTO_SIZE.NONE
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # ミニ見出し（収束の宣言）
    set_text(ctf.paragraphs[0], "頑健な結論", Pt(11), SUB, weight="medium")
    ph = ctf.add_paragraph()
    ph.space_before = Pt(6)
    add_rich_runs(ph, conclusion.get("headline", ""), base_size=Pt(19),
                  base_color=INK, bold_color=INK, line_spacing=1.2, weight="bold")
    if conclusion.get("detail"):
        pdt = ctf.add_paragraph()
        pdt.space_before = Pt(8)
        add_rich_runs(pdt, conclusion["detail"], base_size=Pt(13),
                      base_color=SUB, bold_color=INK, line_spacing=1.4)

    # 各手法行→結論箱への細い矢印（回転矩形。コネクタ禁止）
    concl_mid_y = area_top + area_h / 2
    arr_x = left_x + method_w + 0.06
    for cyc in row_centers:
        _add_line(slide, arr_x, cyc, concl_x - 0.04, concl_mid_y, ACCENT, weight_pt=1.2)
    # 収束先を示す右向き矢印（結論箱直前）
    head = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(concl_x - arrow_w - 0.02),
        Inches(concl_mid_y - 0.13), Inches(arrow_w), Inches(0.26)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = ACCENT
    head.line.fill.background()
    _shadow_off(head)

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 優先度別アクションスライド（優先度バッジ＋期間ピル）
def add_priority_actions_slide(prs, title, actions, blank,
                               eyebrow=None, source=None, page_num=None):
    """優先度別アクション — 各行に [優先度バッジ 高/中/低] ＋ 見出し ＋ 詳細 ＋ [期間ピル]。

    白カード行 + 優先度色バッジ（ピル）+ 期間ピルで行単位で読み取りやすくする。
    （優先度＋期間＋詳細だけで足りる場合は `add_recommendation_slide` でも可。）

    Args:
        actions: list of `{"priority":"高"/"中"/"低", "title":"…", "detail":"…1文",
                          "timeframe":"短期・1年以内"}`
        eyebrow: タイトル直上のアイブロウ
    優先度色: 高=RED_ACCENT、中=CHIP["amber"]、低=MUTED。
    """
    PRIORITY_COLORS = {"高": RED_ACCENT, "中": CHIP["amber"], "低": MUTED}
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title, eyebrow=eyebrow)
    content_y = sub_y + 0.10

    n = len(actions)
    available_h = 6.45 - content_y
    gap = 0.14
    card_h = min(1.35, (available_h - gap * (n - 1)) / max(n, 1))

    badge_w = 0.95
    badge_x = 0.5
    text_x = badge_x + badge_w + 0.30
    pill_w = 2.40
    pill_x = 12.83 - pill_w
    text_w = pill_x - text_x - 0.25

    for i, act in enumerate(actions):
        y = content_y + i * (card_h + gap)
        p_color = PRIORITY_COLORS.get(act.get("priority", "中"), CHIP["amber"])

        # カード地（白カード）
        _card(slide, 0.5, y, 12.33, card_h)

        # 優先度バッジ（色付きピル・白文字）
        badge = _rect(slide, badge_x + 0.12, y + card_h / 2 - 0.26,
                      badge_w, 0.52, p_color, round_=True, radius_in=0.08)
        btf = badge.text_frame
        btf.word_wrap = True
        btf.margin_top = Emu(0)
        btf.margin_bottom = Emu(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        set_text(bp, act.get("priority", "中"), Pt(15), WHITE, weight="bold")

        # 見出し（インク・bold）
        tb = slide.shapes.add_textbox(Inches(text_x), Inches(y + 0.10),
                                      Inches(text_w), Inches(card_h - 0.18))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(tf.paragraphs[0], act.get("title", ""), Pt(15), INK, weight="bold")
        if act.get("detail"):
            pdc = tf.add_paragraph()
            pdc.space_before = Pt(3)
            add_rich_runs(pdc, act["detail"], base_size=Pt(12),
                          base_color=SUB, bold_color=INK, line_spacing=1.3)

        # 期間ピル（小さめ・枠線）
        if act.get("timeframe"):
            pill = _rect(slide, pill_x, y + card_h / 2 - 0.20,
                         pill_w - 0.15, 0.40, CARD, p_color, line_w=Emu(12700),
                         round_=True, radius_in=0.08)
            ptf = pill.text_frame
            ptf.word_wrap = True
            ptf.margin_top = Emu(0)
            ptf.margin_bottom = Emu(0)
            ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = ptf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            set_text(pp, act["timeframe"], Pt(11), p_color, weight="medium")

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# アクションアイテムスライド（☐チェックリスト・締め寄り）
def add_action_items_slide(prs, title, items, blank,
                           eyebrow=None, brand_line=None, source=None, page_num=None):
    """アクションアイテム（☐ チェックリスト）— 完結したアクション文を縦に並べる。

    各行は白カード。☐（関数側で付与）＋ アクション文。下部に任意のブランド行を置き、
    結論・締め寄りの体裁にする。

    Args:
        items: list of str（各「完結したアクション文」。☐は関数が付与）
        brand_line: 省略可（例「APOLLO ・ 特許ランドスケープ分析」）。締め寄りのタグライン。
        eyebrow: タイトル直上のアイブロウ
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title, eyebrow=eyebrow)
    content_y = sub_y + 0.10

    brand_h = 0.55 if brand_line else 0.0
    area_bottom = 6.45 - (brand_h + 0.15 if brand_line else 0.0)
    available_h = area_bottom - content_y

    n = len(items)
    gap = 0.14
    row_h = min(1.0, (available_h - gap * (n - 1)) / max(n, 1))
    total_h = row_h * n + gap * (n - 1)
    start_y = content_y + (available_h - total_h) / 2  # 縦中央寄せ

    box_x = 0.5
    box_w = 12.33
    check_w = 0.42

    for i, item in enumerate(items):
        y = start_y + i * (row_h + gap)
        _card(slide, box_x, y, box_w, row_h)

        # ☐ チェックボックス（角丸枠・ACCENT 線）
        chk = _rect(slide, box_x + 0.28, y + row_h / 2 - check_w / 2,
                    check_w, check_w, CARD, ACCENT, line_w=Emu(19050),
                    round_=True, radius_in=0.06)

        # アクション文
        tb = slide.shapes.add_textbox(Inches(box_x + 0.28 + check_w + 0.22),
                                      Inches(y + 0.06),
                                      Inches(box_w - check_w - 0.85), Inches(row_h - 0.12))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_rich_runs(tf.paragraphs[0], item, base_size=Pt(14),
                      base_color=INK, bold_color=INK, line_spacing=1.3)

    # ブランド行（締め寄り・DEEP 濃紺）
    if brand_line:
        bb = _rect(slide, 0.5, area_bottom + 0.15, 12.33, brand_h, DEEP, round_=True)
        btf = bb.text_frame
        btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        set_text(bp, brand_line, Pt(13), WHITE, weight="medium")

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide


# 2x2マトリクススライド（旧式・軸ラベルのみ）
def add_matrix_slide(prs, title, sub_message, quadrants, blank,
                     x_label="→ 成長率", y_label="↑ 規模",
                     source=None, page_num=None):
    """2x2マトリクス（4象限）— ティント地 + カテゴリ色左バー

    quadrants: {"TL":{"title":"新興","items":["A社"]}, "TR":..., "BL":..., "BR":...}
    """
    slide = prs.slides.add_slide(blank)
    _set_bg(slide)
    sub_y = add_title_shape(slide, title)
    content_y = add_sub_message(slide, sub_message, y=sub_y)

    mx = 1.5
    my = content_y + 0.2
    mw = 5.5
    mh = 6.0 - content_y
    half_w = mw / 2
    half_h = mh / 2
    quad_keys = {"TL": "blue", "TR": "teal", "BL": "violet", "BR": "amber"}
    positions = {"TL": (mx, my), "TR": (mx + half_w, my),
                 "BL": (mx, my + half_h), "BR": (mx + half_w, my + half_h)}

    for key, pos in positions.items():
        q = quadrants.get(key, {})
        ck = quad_keys[key]
        qx, qy = pos

        _rect(slide, qx, qy, half_w - 0.05, half_h - 0.05, TINT[ck], round_=True)
        _bar_left(slide, qx, qy, half_h - 0.05, CAT[ck])

        txBox = slide.shapes.add_textbox(
            Inches(qx + 0.18), Inches(qy + 0.1),
            Inches(half_w - 0.4), Inches(0.35)
        )
        set_text(txBox.text_frame.paragraphs[0], q.get("title", ""), Pt(13), INK, weight="semibold")

        items = q.get("items", [])
        if items:
            txBox2 = slide.shapes.add_textbox(
                Inches(qx + 0.22), Inches(qy + 0.5),
                Inches(half_w - 0.45), Inches(half_h - 0.7)
            )
            tf = txBox2.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items[:5]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                set_text(p, f"・{item}", Pt(10), SUB)

    # 軸ラベル（折返し禁止）
    txX = slide.shapes.add_textbox(Inches(mx + mw/2 - 0.75), Inches(my + mh + 0.05),
                                   Inches(1.5), Inches(0.25))
    txX.text_frame.word_wrap = False
    set_text(txX.text_frame.paragraphs[0], x_label, Pt(10), MUTED, weight="medium")
    txY = slide.shapes.add_textbox(Inches(mx - 1.15), Inches(my + mh/2 - 0.15),
                                   Inches(1.05), Inches(0.3))
    txY.text_frame.word_wrap = False
    txY.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_text(txY.text_frame.paragraphs[0], y_label, Pt(10), MUTED, weight="medium")

    if source:
        add_source_label(slide, source)
    add_bottom_bar_and_footer(slide, page_num)
    return slide
